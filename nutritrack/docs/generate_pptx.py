"""
Generate NutriTrack Certification Defense PowerPoint Presentation.

Oral Defense Structure (90 minutes):
  - 60 min: E1-E5 presentation + live demo + E3 kickoff simulation
  - 10 min: Q&A with jury on E6
  - 10 min: E7 presentation
  - 10 min: Jury questions

Slide decks:
  E1 (5 min)  ~5 slides
  E2 (20 min) ~18 slides
  E3 (10 min) ~8 slides
  E4 (15 min) ~14 slides
  E5 (10 min) ~10 slides
  E6 (10 min Q&A) ~8 slides
  E7 (10 min) ~10 slides
  Total: ~73+ slides
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Colours ──────────────────────────────────────────────────────────
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x3A, 0x86, 0xFF)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
ACCENT_ORANGE = RGBColor(0xF3, 0x97, 0x1E)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
MEDIUM_GRAY = RGBColor(0x88, 0x88, 0x88)
DARK_TEXT = RGBColor(0x2C, 0x2C, 0x2C)
BLOCK1_COLOR = RGBColor(0x3A, 0x86, 0xFF)
BLOCK2_COLOR = RGBColor(0x2E, 0xCC, 0x71)
BLOCK3_COLOR = RGBColor(0xF3, 0x97, 0x1E)
BLOCK4_COLOR = RGBColor(0x9B, 0x59, 0xB6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helpers ──────────────────────────────────────────────────────────


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=18,
    bold=False,
    color=WHITE,
    align=PP_ALIGN.LEFT,
    font_name="Calibri",
):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_bullet_slide(
    slide, left, top, width, height, bullets, font_size=16, color=WHITE, level_colors=None, font_name="Calibri"
):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, level) in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size - level * 2)
        p.font.name = font_name
        p.level = level
        if level_colors and level < len(level_colors):
            p.font.color.rgb = level_colors[level]
        else:
            p.font.color.rgb = color
        p.space_after = Pt(4)
    return txBox


def add_section_divider(title, subtitle, accent_color, eval_label, timing):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BG)
    # Accent bar left
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent_color
    shape.line.fill.background()
    # Eval label
    add_textbox(slide, 0.6, 1.0, 4, 0.6, eval_label, 20, False, accent_color)
    # Title
    add_textbox(slide, 0.6, 1.8, 11, 1.2, title, 44, True, WHITE)
    # Subtitle
    add_textbox(slide, 0.6, 3.2, 10, 1.0, subtitle, 22, False, LIGHT_GRAY)
    # Timing
    add_textbox(slide, 0.6, 5.5, 4, 0.5, timing, 16, False, MEDIUM_GRAY)
    return slide


def add_content_slide(title, accent_color=ACCENT_BLUE, notes_text=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    # Top accent line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent_color
    shape.line.fill.background()
    # Title
    add_textbox(slide, 0.5, 0.25, 12, 0.7, title, 28, True, WHITE)
    # Divider line under title
    shape2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(2.5), Inches(0.04))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = accent_color
    shape2.line.fill.background()
    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text
    return slide


def add_table_to_slide(slide, left, top, width, height, headers, rows, header_color=ACCENT_BLUE, font_size=12):
    cols = len(headers)
    row_count = len(rows) + 1
    table_shape = slide.shapes.add_table(row_count, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table
    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(font_size)
            para.font.bold = True
            para.font.color.rgb = WHITE
            para.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
    # Data rows
    for i, row_data in enumerate(rows):
        for j, val in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(font_size)
                para.font.color.rgb = DARK_TEXT
                para.font.name = "Calibri"
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return table


# =====================================================================
#  TITLE SLIDE
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_textbox(slide, 1, 1.0, 11, 1.5, "NutriTrack", 60, True, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 1, 2.5, 11, 0.8, "Fitness Nutrition Tracking Platform", 30, False, ACCENT_BLUE, PP_ALIGN.CENTER)
add_textbox(
    slide,
    1,
    3.8,
    11,
    0.6,
    "Expert en infrastructures de donnees massives (Data Engineer)",
    20,
    False,
    LIGHT_GRAY,
    PP_ALIGN.CENTER,
)
add_textbox(slide, 1, 4.4, 11, 0.5, "RNCP37638 - Level 7 (Master's equivalent)", 18, False, LIGHT_GRAY, PP_ALIGN.CENTER)
add_textbox(slide, 1, 5.5, 11, 0.5, "[Candidate Name]  |  March 2026", 18, False, MEDIUM_GRAY, PP_ALIGN.CENTER)

# =====================================================================
#  AGENDA SLIDE
# =====================================================================
slide = add_content_slide("Defense Agenda - 90 Minutes", ACCENT_BLUE)
add_table_to_slide(
    slide,
    0.8,
    1.4,
    11.5,
    4.5,
    ["Section", "Evaluation", "Competencies", "Duration"],
    [
        ["Block 1: Steer a Data Project", "E1", "C1", "5 min"],
        ["Block 1: Steer a Data Project", "E2", "C1, C2, C3, C4, C6", "20 min"],
        ["Block 1: Kickoff Meeting Simulation", "E3", "C5, C6, C7", "10 min"],
        ["Block 2: Data Collection & API", "E4", "C8, C9, C10, C11, C12", "15 min"],
        ["Block 3: Data Warehouse", "E5", "C13, C14, C15", "10 min"],
        ["LIVE DEMO", "--", "C8-C15", "~ included above"],
        ["Block 3: DW Maintenance (Q&A)", "E6", "C16, C17", "10 min"],
        ["Block 4: Data Lake", "E7", "C18, C19, C20, C21", "10 min"],
        ["Jury Questions", "--", "All", "10 min"],
    ],
    ACCENT_BLUE,
    13,
)

# =====================================================================
#  TECH STACK OVERVIEW
# =====================================================================
slide = add_content_slide("Technology Stack Overview", ACCENT_BLUE)
add_table_to_slide(
    slide,
    0.8,
    1.4,
    11.5,
    4.5,
    ["Layer", "Technology", "Purpose"],
    [
        ["Infrastructure", "Docker Compose (10 services)", "Container orchestration"],
        ["Database", "PostgreSQL 16-alpine", "OLTP (app schema) + OLAP (dw schema)"],
        ["Cache / Queue", "Redis 7-alpine", "Session cache + Celery broker"],
        ["Object Storage", "MinIO (S3-compatible)", "Data Lake: bronze / silver / gold"],
        ["ETL Orchestration", "Apache Airflow 2.8.1", "6 DAGs, CeleryExecutor"],
        ["Big Data Analytics", "DuckDB", "In-memory Parquet queries (3M+ products)"],
        ["REST API", "FastAPI (async)", "JWT auth, OpenAPI docs, 10+ endpoints"],
        ["ORM", "SQLAlchemy 2.0 (async)", "Async PostgreSQL access"],
        ["Frontend", "Streamlit + Plotly", "Interactive nutrition dashboards"],
    ],
    ACCENT_BLUE,
    13,
)

# =====================================================================
#  ARCHITECTURE SLIDE
# =====================================================================
slide = add_content_slide("System Architecture", ACCENT_BLUE)
arch_text = (
    "                              Streamlit UI (:8501)\n"
    "                                    |\n"
    "                              FastAPI REST (:8000)\n"
    "                             /      |      \\\n"
    "                PostgreSQL 16   Redis 7    MinIO\n"
    "                 (:5432)       (:6379)    (:9000)\n"
    "                app | dw                bronze|silver|gold\n"
    "                    |                        |\n"
    "              Airflow 2.8 ───────────────────+\n"
    "              Webserver (:8080)\n"
    "              Scheduler + Worker (Celery)"
)
add_textbox(slide, 1.5, 1.5, 10, 5.0, arch_text, 18, False, ACCENT_GREEN, font_name="Courier New")
add_textbox(
    slide,
    0.8,
    6.3,
    11,
    0.5,
    "All services containerized in docker-compose.yml (211 lines) with healthchecks and volumes",
    14,
    False,
    LIGHT_GRAY,
)

# =====================================================================
#  E1 - INTERVIEW GRIDS (5 min, ~5 slides)
# =====================================================================
add_section_divider(
    "Need Analysis & Interview Grids",
    "C1: Analyze a data project need expression in a feasibility study",
    BLOCK1_COLOR,
    "E1 - Case Study",
    "5 minutes | Block 1",
)

# E1 - Context
slide = add_content_slide("E1 | Project Context & Business Need", BLOCK1_COLOR)
add_bullet_slide(
    slide,
    0.8,
    1.3,
    11,
    5.5,
    [
        ("Business Need", 0),
        ("Users want to track food intake and receive healthier product alternatives", 1),
        ("Nutritionists need aggregated dietary data for research", 1),
        ("Administrators need data quality monitoring and RGPD compliance", 1),
        ("Data Source: Open Food Facts", 0),
        ("3M+ food products, open-source collaborative database", 1),
        ("Nutri-Score (A-E), NOVA groups (1-4), full nutritional values", 1),
        ("Available via REST API, Parquet export, and web scraping", 1),
        ("Feasibility Validated", 0),
        ("No proprietary data licensing needed (open data)", 1),
        ("SMART objectives defined for 3-month delivery", 1),
    ],
    color=WHITE,
    level_colors=[ACCENT_BLUE, LIGHT_GRAY],
)

# E1 - Grid: Data Producers
slide = add_content_slide("E1 | Interview Grid: Data Producers", BLOCK1_COLOR)
add_table_to_slide(
    slide,
    0.5,
    1.3,
    12,
    4.8,
    ["#", "Question", "Domain"],
    [
        ["1", "What food product databases do you maintain or contribute to?", "Business activity"],
        ["2", "How frequently is product information updated?", "Data characterization"],
        ["3", "What metadata accompanies each product entry?", "Metadata"],
        ["4", "In what format is data stored/exported (JSON, CSV, Parquet)?", "Storage format"],
        ["5", "What access methods are available (REST API, bulk export)?", "Data access"],
        ["6", "Are there rate limits or authentication requirements?", "Access constraints"],
        ["7", "What data quality issues have been observed?", "Data quality"],
        ["8", "How are user contributions validated?", "Data governance"],
        ["9", "Is there a data retention or archival policy?", "Data lifecycle"],
        ["10", "Are there personal data elements in the product data?", "RGPD"],
    ],
    BLOCK1_COLOR,
    11,
)

# E1 - Grid: Data Consumers
slide = add_content_slide("E1 | Interview Grid: Data Consumers", BLOCK1_COLOR)
add_table_to_slide(
    slide,
    0.5,
    1.3,
    12,
    4.8,
    ["#", "Question", "Domain"],
    [
        ["1", "What nutritional analyses do you perform today?", "Business need"],
        ["2", "How do you currently search for product information?", "Usage patterns"],
        ["3", "What level of nutritional detail do you need?", "Data granularity"],
        ["4", "Do you need historical data for product reformulations?", "Temporal analysis"],
        ["5", "What aggregation level is useful (product/brand/category)?", "Analytical needs"],
        ["6", "How should data be delivered (API, dashboard, export)?", "Data access"],
        ["7", "What response time is acceptable for searches?", "Performance SLA"],
        ["8", "Do you need individual user meal tracking (RGPD)?", "Consent requirements"],
        ["9", "What accessibility accommodations are needed?", "Accessibility"],
        ["10", "Regulatory constraints on nutritional data display?", "Regulatory"],
    ],
    BLOCK1_COLOR,
    11,
)

# E1 - Synthesis
slide = add_content_slide("E1 | Synthesis Note Summary", BLOCK1_COLOR)
add_bullet_slide(
    slide,
    0.8,
    1.3,
    11,
    5.5,
    [
        ("SMART Objectives", 0),
        ("S: Build a nutrition platform integrating 5 data source types", 1),
        ("M: Process 100,000+ products with 95%+ completeness on key fields", 1),
        ("A: Open-source stack within Docker Compose environment", 1),
        ("R: Aligned with Open Food Facts model + EU nutritional standards", 1),
        ("T: Deliver within a 3-month sprint cycle", 1),
        ("Pre-Project Includes", 0),
        ("Macro technical recommendations (PostgreSQL, Airflow, MinIO)", 1),
        ("RGPD compliance: consent tracking, data retention, pseudonymization", 1),
        ("Accessibility: web-based UI, responsive Streamlit layout", 1),
    ],
    color=WHITE,
    level_colors=[ACCENT_BLUE, LIGHT_GRAY],
)

# =====================================================================
#  E2 - PROFESSIONAL SITUATION (20 min, ~18 slides)
# =====================================================================
add_section_divider(
    "Data Topography, Architecture & Monitoring",
    "C1, C2, C3, C4, C6: Need analysis, data mapping, technical framework, monitoring, supervision",
    BLOCK1_COLOR,
    "E2 - Professional Situation",
    "20 minutes | Block 1",
)

# E2 - Data Topography: Semantics
slide = add_content_slide("E2 | C2 - Data Topography: Business Glossary", BLOCK1_COLOR)
add_table_to_slide(
    slide,
    0.5,
    1.3,
    12,
    5.0,
    ["Business Object", "Description", "Key Metadata"],
    [
        [
            "Product",
            "Food product identified by barcode",
            "barcode, name, brand, category, nutrition/100g, Nutri-Score, NOVA",
        ],
        ["Brand", "Product manufacturer", "brand_name, parent_company"],
        ["Category", "Hierarchical food classification", "category_name, parent_category, level"],
        ["User", "App user tracking nutrition", "user_id (UUID), email, role, dietary_goal, consent"],
        ["Meal", "Logged eating occasion", "meal_type, meal_date, user_id"],
        ["Meal Item", "Single product within a meal", "product_id, quantity_g, computed nutrition"],
        ["Nutri-Score", "Nutritional quality grade", "A (best) to E (worst), score -15 to +40"],
        ["NOVA Group", "Processing classification", "1 (minimal) to 4 (ultra-processed)"],
    ],
    BLOCK1_COLOR,
    12,
)

# E2 - Data Models
slide = add_content_slide("E2 | C2 - Data Models", BLOCK1_COLOR)
add_table_to_slide(
    slide,
    0.8,
    1.3,
    11.5,
    2.5,
    ["Data Type", "Model", "Storage"],
    [
        ["Structured", "Relational (PostgreSQL)", "Normalized tables with FK constraints"],
        ["Semi-structured", "JSON", "API responses, scraping results, catalog metadata"],
        ["Unstructured", "Columnar / Binary", "Parquet files (3M+ products), compressed CSV"],
    ],
    BLOCK1_COLOR,
    14,
)
add_textbox(slide, 0.8, 4.2, 11, 0.5, "Two schemas in PostgreSQL:", 18, True, ACCENT_BLUE)
add_bullet_slide(
    slide,
    0.8,
    4.8,
    11,
    2.0,
    [
        ("app schema: Operational (OLTP) - products, users, meals, RGPD registry", 0),
        ("dw schema: Analytical (OLAP) - star schema with dimensions & fact tables", 0),
    ],
    color=LIGHT_GRAY,
    font_size=16,
)

# E2 - Flux Matrix
slide = add_content_slide("E2 | C2 - Treatments & Data Flows (Flux Matrix)", BLOCK1_COLOR)
add_table_to_slide(
    slide,
    0.3,
    1.3,
    12.5,
    5.0,
    ["Source", "Format", "Script / Treatment", "Target", "Frequency"],
    [
        ["OFF REST API", "JSON", "extract_off_api.py", "data/raw/api/", "Daily"],
        ["OFF Parquet Export", "Parquet", "extract_off_parquet.py (DuckDB)", "data/raw/parquet/", "Weekly"],
        ["ANSES/EFSA Websites", "HTML", "extract_scraping.py (BS4)", "data/raw/scraping/", "Weekly"],
        ["PostgreSQL (operational)", "SQL", "extract_from_db.py", "data/raw/database/", "On-demand"],
        ["DuckDB Analytics", "Parquet", "extract_duckdb.py", "data/raw/duckdb/", "Weekly"],
        ["All raw sources", "Mixed", "aggregate_clean.py", "data/cleaned/", "Daily"],
        ["Cleaned dataset", "Parquet", "import_to_db.py", "app schema (PostgreSQL)", "Daily"],
        ["Operational DB", "SQL", "etl_load_warehouse.py (Airflow)", "dw schema (PostgreSQL)", "Daily"],
        ["All layers", "Mixed", "etl_datalake_ingest.py", "MinIO (bronze/silver/gold)", "Daily"],
    ],
    BLOCK1_COLOR,
    11,
)

# E2 - Data Access Conditions
slide = add_content_slide("E2 | C2 - Data Access & Usage Conditions", BLOCK1_COLOR)
add_table_to_slide(
    slide,
    0.5,
    1.3,
    12,
    4.0,
    ["Role", "Operational DB (app)", "Data Warehouse (dw)", "Data Lake (MinIO)"],
    [
        ["app_readonly", "SELECT: products, brands, categories", "No access", "No access"],
        ["nutritionist_role", "SELECT all except users", "SELECT datamart views", "Read silver, gold"],
        ["admin_role", "Full DML access", "Full access", "Full access all buckets"],
        ["etl_service", "INSERT, UPDATE", "Full DML for ETL", "Write all data buckets"],
        ["nutritrack (app)", "Full operational access", "Full access", "N/A"],
    ],
    BLOCK1_COLOR,
    12,
)

# E2 - C3 Architecture Decisions
slide = add_content_slide("E2 | C3 - Architecture Decisions", BLOCK1_COLOR)
add_table_to_slide(
    slide,
    0.5,
    1.3,
    12,
    5.0,
    ["Decision", "Chosen", "Alternative", "Rationale"],
    [
        [
            "DB for OLTP/OLAP",
            "Single PostgreSQL (separate schemas)",
            "Separate instances",
            "Simpler deployment; adequate for current scale",
        ],
        ["ETL orchestrator", "Apache Airflow", "Prefect, Luigi", "Industry standard, web UI, Celery parallelism"],
        ["Data lake storage", "MinIO (self-hosted)", "AWS S3", "Open-source, same S3 API, no cloud costs"],
        ["API framework", "FastAPI (async)", "Flask, Django", "Native async, auto OpenAPI docs, Pydantic"],
        ["Big data analytics", "DuckDB", "Apache Spark", "No cluster needed, SQL interface, reads Parquet"],
        ["DW approach", "Bottom-up (Kimball)", "Top-down (Inmon)", "Faster value delivery, incremental datamarts"],
    ],
    BLOCK1_COLOR,
    12,
)

# E2 - C3 Non-Functional Requirements
slide = add_content_slide("E2 | C3 - Non-Functional Requirements", BLOCK1_COLOR)
add_table_to_slide(
    slide,
    0.8,
    1.3,
    11.5,
    3.5,
    ["Constraint", "Requirement", "Solution"],
    [
        ["Performance", "Product search < 1s", "GIN index for full-text search, connection pooling"],
        ["Scalability", "Handle 3M+ products", "DuckDB for in-memory Parquet analytics"],
        ["Reliability", "Zero data loss on failures", "Airflow retries (2-3x), extraction audit log"],
        ["Security", "No plaintext passwords", "bcrypt hashing, JWT HS256, role-based access"],
        ["Availability", "Health monitoring", "Docker healthchecks, /health endpoints, Airflow UI"],
    ],
    BLOCK1_COLOR,
    13,
)

# E2 - C3 RGPD
slide = add_content_slide("E2 | C3 - RGPD Compliance Processes", BLOCK1_COLOR)
add_table_to_slide(
    slide,
    0.5,
    1.3,
    12,
    3.5,
    ["Data Category", "Legal Basis", "Retention", "Security Measures"],
    [
        [
            "User account (email, username, hash)",
            "Consent Art. 6.1.a",
            "3 years after last login",
            "bcrypt hashing, encrypted storage",
        ],
        [
            "User profile (age, activity, goal)",
            "Consent Art. 6.1.a",
            "3 years or consent withdrawal",
            "SHA256 pseudonymized in analytics",
        ],
        [
            "Meal tracking (logs, history)",
            "Consent Art. 6.1.a",
            "2 years after creation",
            "UUID association, deletable on request",
        ],
        [
            "Product data (Open Food Facts)",
            "Legitimate interest Art. 6.1.f",
            "Indefinite (public data)",
            "Public data, no personal info",
        ],
    ],
    BLOCK1_COLOR,
    12,
)
add_bullet_slide(
    slide,
    0.8,
    5.2,
    11,
    2.0,
    [
        ("Automated cleanup: rgpd_cleanup_expired_data() - deletes old meals, deactivates expired users", 0),
        ("Data lake: no personal data stored (only product data + anonymized aggregations)", 0),
        ("Lifecycle rules: bronze auto-expires at 90 days, backups at 30 days", 0),
    ],
    color=LIGHT_GRAY,
    font_size=14,
)

# E2 - C3 Eco-Responsibility
slide = add_content_slide("E2 | C3 - Eco-Responsibility Strategy (RGESN)", BLOCK1_COLOR)
add_bullet_slide(
    slide,
    0.8,
    1.3,
    11,
    5.5,
    [
        ("Following the French RGESN framework:", 0),
        ("Efficient queries: DuckDB processes 3M+ products in-memory; composite indexes reduce scan volume", 0),
        ("Minimal data transfers: API pagination limits response size; only delta changes in DW ETL", 0),
        ("Data lifecycle management: bronze auto-expires 90 days; backups purged at 30 days", 0),
        ("Lightweight containers: Alpine-based images (PostgreSQL, Redis) minimize resource usage", 0),
        ("Batch processing: all ETL runs during off-peak hours (02:00-06:00 UTC)", 0),
    ],
    color=LIGHT_GRAY,
    font_size=16,
)

# E2 - C4 Monitoring
slide = add_content_slide("E2 | C4 - Technical & Regulatory Monitoring", BLOCK1_COLOR)
add_bullet_slide(
    slide,
    0.8,
    1.3,
    11,
    5.5,
    [
        ("Monitoring Theme: Apache Airflow - ETL Orchestration", 0),
        ("Schedule: Minimum 1h/week reviewing DAG execution, pipeline health, data quality", 1),
        ("Aggregation Tools", 0),
        ("Airflow Web UI: DAG run history, task duration, success/failure rates, Gantt charts", 1),
        ("extraction_log table: audit trail (records extracted/loaded/rejected per source)", 1),
        ("cleaning_report.json: per-column completeness statistics", 1),
        ("Source Reliability Criteria", 0),
        ("Open Food Facts: 30,000+ contributors, validated before publication", 1),
        ("EU Regulation 1169/2011: official EU legislative source", 1),
        ("ANSES: French national health authority, peer-reviewed values", 1),
        ("Alert System", 0),
        ("Airflow email_on_failure: True -> admin@nutritrack.local", 1),
    ],
    color=WHITE,
    level_colors=[ACCENT_BLUE, LIGHT_GRAY],
)

# E2 - C6 Supervision
slide = add_content_slide("E2 | C6 - Project Supervision & Tracking", BLOCK1_COLOR)
add_bullet_slide(
    slide,
    0.8,
    1.3,
    11,
    5.5,
    [
        ("Exchange Facilitation", 0),
        ("Makefile: standardized commands (make pipeline, make backup, make test-api)", 1),
        ("Tracking Tools", 0),
        ("Airflow UI (port 8080): real-time DAG status, task logs, historical data", 1),
        ("extraction_log table: every data operation tracked with timestamps and counts", 1),
        ("Documented Rituals", 0),
        ("Daily ETL: sequential schedule (02:00-06:00 UTC) with dependency chains", 1),
        ("Weekly data quality review using cleaning report", 1),
        ("Indicators Updated Throughout Project", 0),
        ("Records extracted/loaded/rejected per source", 1),
        ("Data completeness % per column", 1),
        ("DAG success rate over time", 1),
    ],
    color=WHITE,
    level_colors=[ACCENT_BLUE, LIGHT_GRAY],
)

# E2 - Lessons Learned
slide = add_content_slide("E2 | Lessons Learned & Trade-offs", BLOCK1_COLOR)
add_bullet_slide(
    slide,
    0.8,
    1.3,
    11,
    5.5,
    [
        ("Key Trade-offs Made", 0),
        ("Single PostgreSQL for OLTP+OLAP: simplicity vs. optimal separation", 1),
        ("MinIO instead of AWS S3: self-hosting + cost savings vs. managed service", 1),
        ("DuckDB vs. Spark: simplicity + no cluster vs. distributed computing", 1),
        ("Bottom-up DW: faster delivery vs. comprehensive top-down design", 1),
        ("Problems Encountered", 0),
        ("Airflow CeleryExecutor setup: Redis broker + PostgreSQL backend + init containers", 1),
        ("NULL handling in SCD Type 2: required IS DISTINCT FROM instead of <> operator", 1),
        ("OFF API rate limiting: required 0.6s delay between requests", 1),
    ],
    color=WHITE,
    level_colors=[ACCENT_ORANGE, LIGHT_GRAY],
)

# =====================================================================
#  E3 - KICKOFF MEETING SIMULATION (10 min, ~8 slides)
# =====================================================================
add_section_divider(
    "Project Kickoff Meeting Simulation",
    "C5, C6, C7: Planning, supervision, communication",
    BLOCK1_COLOR,
    "E3 - Role Play / Kickoff Meeting",
    "10 minutes | Block 1",
)

# E3 - Pre-project
slide = add_content_slide("E3 | Pre-Project Summary", BLOCK1_COLOR)
add_bullet_slide(
    slide,
    0.8,
    1.3,
    11,
    5.5,
    [
        ("Project: NutriTrack - Fitness Nutrition Tracking Platform", 0),
        ("Objective: End-to-end data engineering platform demonstrating all 21 competencies", 0),
        ("Scope: Data collection (5 sources) -> Cleaning -> DB -> DW -> Data Lake -> API -> UI", 0),
        ("Team Composition", 0),
        ("Data Engineer (lead): Python, SQL, Docker, Airflow, ETL design", 1),
        ("Backend Developer: FastAPI, SQLAlchemy, JWT authentication", 1),
        ("Frontend Developer: Streamlit, Plotly, UX design", 1),
        ("DBA / Data Architect: PostgreSQL, star schema, SCD, performance tuning", 1),
        ("DevOps: Docker Compose, MinIO, CI/CD", 1),
    ],
    color=WHITE,
    level_colors=[ACCENT_BLUE, LIGHT_GRAY],
)

# E3 - Roadmap
slide = add_content_slide("E3 | C5 - Project Roadmap (6 Phases)", BLOCK1_COLOR)
add_table_to_slide(
    slide,
    0.5,
    1.3,
    12,
    4.5,
    ["Phase", "Duration", "Deliverables", "Milestone"],
    [
        ["1. Planning & Design", "2 weeks", "Interview grids, architecture study", "Architecture validated"],
        ["2. Data Collection", "3 weeks", "5 extraction scripts, SQL queries", "Data pipeline operational"],
        ["3. Database & API", "2 weeks", "Operational DB, REST API, import scripts", "API functional + docs"],
        ["4. Data Warehouse", "2 weeks", "Star schema, ETL DAGs, SCD procedures", "DW populated"],
        ["5. Data Lake", "2 weeks", "MinIO, medallion architecture, catalog", "Full lake operational"],
        ["6. Frontend & Integration", "1 week", "Streamlit app, testing, demo prep", "Platform demo-ready"],
    ],
    BLOCK1_COLOR,
    12,
)

# E3 - Calendar
slide = add_content_slide("E3 | C5 - Production Calendar with Effort Weighting", BLOCK1_COLOR)
add_table_to_slide(
    slide,
    0.3,
    1.3,
    12.5,
    5.0,
    ["Week", "Tasks", "Deliverable", "Story Points"],
    [
        ["W1-W2", "Interview grids, data topography, architecture", "E1+E2 deliverables", "21"],
        ["W3-W4", "REST API extraction, Parquet, scraping, DuckDB", "5 extraction scripts", "16"],
        ["W5-W6", "SQL query optimization, import script, DB schema", "7 queries + MPD", "21"],
        ["W7", "REST API development, JWT auth, endpoints", "API + OpenAPI docs", "13"],
        ["W8-W9", "Star schema, warehouse ETL DAGs, SCD procedures", "DW + 6 DAGs", "21"],
        ["W10-W11", "MinIO setup, medallion, catalog, governance", "Data lake", "16"],
        ["W12", "Streamlit frontend, integration testing, demo", "Working platform", "13"],
    ],
    BLOCK1_COLOR,
    12,
)
add_textbox(
    slide,
    0.8,
    6.5,
    11,
    0.5,
    "Effort weighting: Fibonacci story points (team estimation, similar to poker planning)",
    14,
    False,
    MEDIUM_GRAY,
)

# E3 - Tracking & Rituals
slide = add_content_slide("E3 | C6 - Tracking Method & Rituals", BLOCK1_COLOR)
add_bullet_slide(
    slide,
    0.8,
    1.3,
    5,
    5.5,
    [
        ("Methodology: Agile (2-week sprints)", 0),
        ("Tracking Tools", 0),
        ("Git: version control for all code + SQL", 1),
        ("Airflow UI: DAG execution monitoring", 1),
        ("Makefile: standardized dev commands", 1),
        ("Docker Compose: reproducible environment", 1),
    ],
    color=WHITE,
    level_colors=[ACCENT_BLUE, LIGHT_GRAY],
    font_size=15,
)
add_bullet_slide(
    slide,
    6.5,
    1.3,
    6,
    5.5,
    [
        ("Sprint Rituals", 0),
        ("Sprint planning (biweekly): goals + backlog", 1),
        ("Daily standup (15 min): blockers + progress", 1),
        ("Sprint review: demo completed features", 1),
        ("Sprint retrospective: process improvements", 1),
    ],
    color=WHITE,
    level_colors=[ACCENT_BLUE, LIGHT_GRAY],
    font_size=15,
)

# E3 - Communication Strategy
slide = add_content_slide("E3 | C7 - Communication Strategy", BLOCK1_COLOR)
add_table_to_slide(
    slide,
    0.5,
    1.3,
    12,
    4.5,
    ["Communication", "Audience", "Format", "Frequency"],
    [
        ["Sprint kickoff", "All stakeholders", "Presentation", "Every 2 weeks"],
        ["Progress update", "Project sponsor", "Email summary", "Weekly"],
        ["Technical demo", "Dev team + nutritionists", "Live demo", "End of each phase"],
        ["API documentation", "Developers", "OpenAPI / Swagger UI", "Continuous"],
        ["User documentation", "End users", "Help pages in Streamlit", "At delivery"],
        ["Final delivery", "Jury + stakeholders", "Oral presentation + report", "End of project"],
    ],
    BLOCK1_COLOR,
    13,
)
add_bullet_slide(
    slide,
    0.8,
    5.5,
    11,
    1.5,
    [
        ("User documentation planned in W12, assigned to frontend developer", 0),
        ("End-user onboarding session planned as part of final demo", 0),
        ("Post-demo feedback questionnaire integrated into communication strategy", 0),
    ],
    color=LIGHT_GRAY,
    font_size=14,
)

# E3 - Orientations communicated
slide = add_content_slide("E3 | C7 - Orientations & Trade-offs Communicated", BLOCK1_COLOR)
add_bullet_slide(
    slide,
    0.8,
    1.3,
    11,
    5.5,
    [
        ("Technical Choices Communicated to Stakeholders:", 0),
        ("Single PostgreSQL for OLTP+OLAP: explained simplicity benefit vs. future scaling need", 1),
        ("MinIO self-hosted: explained cost savings and S3 API compatibility", 1),
        ("DuckDB for big data: explained no cluster overhead, SQL familiarity", 1),
        ("Bottom-up DW: explained faster value delivery with incremental datamarts", 1),
        ("FastAPI over Flask/Django: explained async performance and auto-documentation", 1),
        ("Communications adapt content to audience:", 0),
        ("Technical team: architecture diagrams, code reviews, SQL optimization notes", 1),
        ("Nutritionists: datamart views, analytical query results", 1),
        ("End users: Streamlit UI walkthrough, search + meal logging demos", 1),
    ],
    color=WHITE,
    level_colors=[ACCENT_BLUE, LIGHT_GRAY],
)

# =====================================================================
#  E4 - DATA COLLECTION, STORAGE & API (15 min, ~14 slides)
# =====================================================================
add_section_divider(
    "Data Collection, Storage & Sharing",
    "C8-C12: Extraction, SQL queries, aggregation, database, REST API",
    BLOCK2_COLOR,
    "E4 - Professional Situation",
    "15 minutes | Block 2",
)

# E4 - C8 Overview
slide = add_content_slide("E4 | C8 - 5 Extraction Scripts Overview", BLOCK2_COLOR)
add_table_to_slide(
    slide,
    0.5,
    1.3,
    12,
    4.5,
    ["Source Type", "Script", "Technology", "Output"],
    [
        ["REST API", "extract_off_api.py", "requests + pagination + rate limiting", "JSON"],
        ["Data File (Big Data)", "extract_off_parquet.py", "DuckDB SQL on 3M+ Parquet", "Parquet/CSV"],
        ["Web Scraping", "extract_scraping.py", "BeautifulSoup + lxml", "JSON/CSV"],
        ["Database", "extract_from_db.py", "SQLAlchemy + PostgreSQL", "Parquet/CSV"],
        ["Big Data Analytics", "extract_duckdb.py", "DuckDB analytical queries", "Parquet"],
    ],
    BLOCK2_COLOR,
    13,
)
add_bullet_slide(
    slide,
    0.8,
    5.2,
    11,
    2.0,
    [
        (
            "All scripts follow C8 structure: entry point, dependency init, external connections, logic rules, error handling, result saving",
            0,
        ),
        ("All scripts versioned on Git", 0),
        ("Each has a corresponding Airflow DAG for automated scheduling", 0),
    ],
    color=LIGHT_GRAY,
    font_size=14,
)

# E4 - C8 REST API extraction
slide = add_content_slide("E4 | C8 - REST API Extraction (extract_off_api.py)", BLOCK2_COLOR)
code = (
    "def search_products(query, page=1, page_size=50):\n"
    "    params = {\n"
    '        "search_terms": query,\n'
    '        "json": 1, "page": page,\n'
    '        "page_size": page_size,\n'
    '        "fields": ",".join(PRODUCT_FIELDS),\n'
    "    }\n"
    '    headers = {"User-Agent": USER_AGENT}\n'
    "    response = requests.get(OFF_API_SEARCH,\n"
    "        params=params, headers=headers, timeout=60)\n"
    "    response.raise_for_status()\n"
    '    return response.json().get("products", [])\n\n'
    "def extract_products_by_category(category, max_pages=10):\n"
    "    for page in range(1, max_pages + 1):\n"
    "        products = search_products(category, page=page)\n"
    "        if not products: break\n"
    "        all_products.extend(products)\n"
    "        time.sleep(0.6)  # Rate limiting"
)
add_textbox(slide, 0.8, 1.3, 11, 5.5, code, 13, False, ACCENT_GREEN, font_name="Courier New")

# E4 - C8 DuckDB extraction
slide = add_content_slide("E4 | C8 - DuckDB Extraction (extract_off_parquet.py)", BLOCK2_COLOR)
code2 = (
    'def extract_with_duckdb(data_path, countries_filter="France"):\n'
    "    con = duckdb.connect()\n"
    '    query = f"""\n'
    "    SELECT code AS barcode, product_name,\n"
    '           "energy-kcal_100g" AS energy_kcal,\n'
    "           nutriscore_grade, nova_group,\n"
    "           completeness AS completeness_score\n"
    "    FROM read_parquet('{data_path}')\n"
    "    WHERE code IS NOT NULL\n"
    "      AND product_name IS NOT NULL\n"
    "      AND LENGTH(code) >= 8\n"
    "      AND countries LIKE '%{countries_filter}%'\n"
    "    ORDER BY completeness DESC\n"
    '    """\n'
    "    df = con.execute(query).fetchdf()  # 3M+ rows in seconds\n"
    "    con.close()\n"
    "    return df"
)
add_textbox(slide, 0.8, 1.3, 11, 5.5, code2, 13, False, ACCENT_GREEN, font_name="Courier New")

# E4 - C8 Web Scraping
slide = add_content_slide("E4 | C8 - Web Scraping (extract_scraping.py)", BLOCK2_COLOR)
add_bullet_slide(
    slide,
    0.8,
    1.3,
    11,
    5.5,
    [
        ("Source: ANSES (French health authority) + EFSA nutritional references", 0),
        ("Technology: BeautifulSoup with lxml parser", 0),
        ("Scraping logic:", 0),
        ("fetch_page(url) -> BeautifulSoup object with User-Agent header", 1),
        ("scrape_anses_guidelines() -> extracts HTML table rows (nutrient, daily_value, unit)", 1),
        ("Graceful fallback to EU Regulation 1169/2011 hardcoded RDA values (25 nutrients)", 1),
        ("Error handling:", 0),
        ("try/except RequestException with timeout=30s", 1),
        ("Fallback data ensures pipeline never fails completely", 1),
        ("Polite crawling with time.sleep(1) between requests", 1),
        ("Output: JSON + CSV with metadata (extraction_date, source, record_count)", 0),
    ],
    color=WHITE,
    level_colors=[ACCENT_GREEN, LIGHT_GRAY],
)

# E4 - C9 SQL Queries
slide = add_content_slide("E4 | C9 - 7 Analytical SQL Queries with Optimizations", BLOCK2_COLOR)
add_table_to_slide(
    slide,
    0.3,
    1.3,
    12.5,
    5.0,
    ["#", "Query Purpose", "Optimization Technique", "Key SQL Feature"],
    [
        ["Q1", "Daily nutrition per user with RDA%", "Composite index (user_id, meal_date)", "JOIN + SUM + GROUP BY"],
        [
            "Q2",
            "Full-text product search + Nutri-Score filter",
            "GIN index on to_tsvector",
            "ts_rank + plainto_tsquery",
        ],
        [
            "Q3",
            "Healthier alternatives in same category",
            "Window function avoids self-join",
            "ROW_NUMBER PARTITION BY",
        ],
        ["Q4", "Weekly trends with moving average", "Window frame (ROWS BETWEEN)", "AVG OVER + LAG"],
        ["Q5", "Nutri-Score distribution (DW)", "Pre-aggregated fact table", "SUM(COUNT(*)) OVER"],
        ["Q6", "Brand quality with NOVA analysis", "CTE + HAVING post-filter", "PERCENTILE_CONT + CASE"],
        ["Q7", "User meal pattern analysis (DW)", "Grouped aggregation", "STDDEV + GROUP BY"],
    ],
    BLOCK2_COLOR,
    12,
)

# E4 - C9 SQL Code example
slide = add_content_slide("E4 | C9 - SQL Query Example: Healthier Alternatives", BLOCK2_COLOR)
sql_code = (
    "WITH product_ranking AS (\n"
    "    SELECT p.product_id, p.barcode,\n"
    "           p.product_name, b.brand_name,\n"
    "           p.nutriscore_grade, p.nutriscore_score,\n"
    "           ROW_NUMBER() OVER (\n"
    "               PARTITION BY p.category_id\n"
    "               ORDER BY p.nutriscore_score ASC\n"
    "           ) AS rank_in_category\n"
    "    FROM app.products p\n"
    "    LEFT JOIN app.brands b ON p.brand_id = b.brand_id\n"
    "    WHERE p.category_id = (\n"
    "        SELECT category_id FROM app.products\n"
    "        WHERE barcode = '3017620422003'\n"
    "    )\n"
    ")\n"
    "SELECT * FROM product_ranking\n"
    "WHERE rank_in_category <= 5;"
)
add_textbox(slide, 0.8, 1.3, 11, 5.5, sql_code, 14, False, ACCENT_GREEN, font_name="Courier New")

# E4 - C10 Cleaning Pipeline
slide = add_content_slide("E4 | C10 - Aggregation & Cleaning Pipeline", BLOCK2_COLOR)
add_bullet_slide(
    slide,
    0.8,
    1.3,
    11,
    5.5,
    [
        ("Pipeline: aggregate_clean.py (466 lines)", 0),
        ("1. Load from all 5 sources (API JSON, Parquet, DuckDB, DB, scraping)", 0),
        ("2. Standardize columns (30+ mappings: 'code' -> 'barcode', 'fat_100g' -> 'fat_g')", 0),
        ("3. Clean barcodes: trim, remove non-digits, validate 8-14 digits", 0),
        ("4. Clean text: normalize whitespace, remove control characters", 0),
        ("5. Remove entries without product_name", 0),
        ("6. Convert kJ -> kcal where kcal is missing", 0),
        ("7. Validate numeric ranges (energy <= 1000, fats <= 100, proteins <= 100)", 0),
        ("8. Normalize Nutri-Score to A-E only; validate NOVA group 1-4", 0),
        ("9. Deduplicate by barcode (keep most complete version)", 0),
        ("Output: products_cleaned.parquet + .csv + cleaning_report.json", 0),
    ],
    color=WHITE,
    level_colors=[ACCENT_GREEN, LIGHT_GRAY],
)

# E4 - C11 Data Models
slide = add_content_slide("E4 | C11 - MERISE Models (MCD / MLD / MPD)", BLOCK2_COLOR)
mcd = (
    "   MCD (Conceptual Model):\n\n"
    "   BRAND ----1:N---- PRODUCT ----N:1---- CATEGORY\n"
    "                        |N                  (self-join)\n"
    "                        |\n"
    "                   MEAL_ITEM\n"
    "                        |N\n"
    "                      MEAL ----N:1---- USER\n\n"
    "   MLD (Logical Model):\n"
    "   products(product_id PK, barcode UK, brand_id FK, category_id FK, ...)\n"
    "   users(user_id PK UUID, email UK, consent_data_processing, ...)\n"
    "   meals(meal_id PK, user_id FK CASCADE, meal_type, meal_date)\n"
    "   meal_items(meal_item_id PK, meal_id FK CASCADE, product_id FK, qty)\n\n"
    "   MPD: sql/init/01_schema_operational.sql (279 lines)\n"
    "   Key: UUID for user_id, GIN index, composite indexes, CHECK constraints"
)
add_textbox(slide, 0.8, 1.3, 11, 5.5, mcd, 14, False, ACCENT_GREEN, font_name="Courier New")

# E4 - C11 RGPD
slide = add_content_slide("E4 | C11 - RGPD Registry & Compliance", BLOCK2_COLOR)
add_bullet_slide(
    slide,
    0.8,
    1.3,
    11,
    5.5,
    [
        ("RGPD Data Registry (app.rgpd_data_registry table):", 0),
        ("4 entries covering: user account, user profile, meal tracking, product data", 1),
        ("Each entry: data_category, legal_basis, retention_period, security_measures", 1),
        ("Automated Cleanup Function:", 0),
        ("rgpd_cleanup_expired_data(): deletes meals > 2 years, deactivates expired users", 1),
        ("Triggered periodically via scheduled maintenance", 1),
        ("Import Script (import_to_db.py):", 0),
        ("Batch upsert: ON CONFLICT (barcode) DO UPDATE", 1),
        ("COALESCE preserves existing data when new record has NULLs", 1),
        ("Batch size: 5,000 records per transaction for optimal throughput", 1),
    ],
    color=WHITE,
    level_colors=[ACCENT_GREEN, LIGHT_GRAY],
)

# E4 - C12 API Endpoints
slide = add_content_slide("E4 | C12 - REST API Endpoints", BLOCK2_COLOR)
add_table_to_slide(
    slide,
    0.3,
    1.3,
    12.5,
    5.0,
    ["Method", "Endpoint", "Auth", "Description"],
    [
        ["POST", "/api/v1/auth/register", "None", "Create user (with RGPD consent)"],
        ["POST", "/api/v1/auth/login", "None", "Authenticate -> JWT token"],
        ["GET", "/api/v1/auth/me", "Bearer", "Current user profile"],
        ["GET", "/api/v1/products/{barcode}", "Bearer", "Get product by barcode"],
        ["GET", "/api/v1/products/?q=...", "Bearer", "Search + filter (nutriscore, nova)"],
        ["GET", "/api/v1/products/{bc}/alternatives", "Bearer", "Healthier alternatives"],
        ["POST", "/api/v1/meals/", "Bearer", "Log a meal with items"],
        ["GET", "/api/v1/meals/", "Bearer", "List user's meals"],
        ["GET", "/api/v1/meals/daily-summary", "Bearer", "Daily nutrition summary"],
        ["GET", "/api/v1/meals/weekly-trends", "Bearer", "Weekly trends + moving avg"],
    ],
    BLOCK2_COLOR,
    11,
)
add_textbox(
    slide,
    0.8,
    6.3,
    11,
    0.5,
    "Auth: JWT HS256 + bcrypt | Docs: /docs (Swagger) + /redoc | RBAC: require_role() decorator",
    14,
    False,
    MEDIUM_GRAY,
)

# E4 - C12 JWT Code
slide = add_content_slide("E4 | C12 - JWT Authentication & RBAC", BLOCK2_COLOR)
jwt_code = (
    "# api/auth/jwt.py\n"
    "pwd_context = CryptContext(schemes=['bcrypt'])\n"
    "oauth2_scheme = OAuth2PasswordBearer(\n"
    "    tokenUrl='/api/v1/auth/login')\n\n"
    "def create_access_token(data: dict):\n"
    "    to_encode = data.copy()\n"
    "    expire = datetime.now(timezone.utc) + timedelta(\n"
    "        minutes=settings.ACCESS_TOKEN_EXPIRE_MINUTES)\n"
    "    to_encode.update({'exp': expire})\n"
    "    return jwt.encode(to_encode,\n"
    "        settings.JWT_SECRET_KEY,\n"
    "        algorithm=settings.JWT_ALGORITHM)\n\n"
    "def require_role(*roles):\n"
    "    async def role_checker(current_user=Depends(...)):\n"
    "        if current_user.role not in roles:\n"
    "            raise HTTPException(403, 'Not authorized')\n"
    "        return current_user\n"
    "    return role_checker"
)
add_textbox(slide, 0.8, 1.3, 11, 5.5, jwt_code, 13, False, ACCENT_GREEN, font_name="Courier New")

# =====================================================================
#  E5 - DATA WAREHOUSE (10 min, ~10 slides)
# =====================================================================
add_section_divider(
    "Data Warehouse",
    "C13, C14, C15: Star schema modeling, DW creation, ETL pipelines",
    BLOCK3_COLOR,
    "E5 - Professional Situation",
    "10 minutes | Block 3",
)

# E5 - C13 Star Schema
slide = add_content_slide("E5 | C13 - Star Schema Design", BLOCK3_COLOR)
schema = (
    "              dim_time          dim_user (SCD2)\n"
    "              time_key PK       user_key PK\n"
    "              full_date         user_hash (SHA256)\n"
    "              day/month/year    age_group, dietary_goal\n"
    "                  |                 |\n"
    "   dim_brand      |    fact_daily_nutrition    dim_category\n"
    "   brand_key PK --+--> user_key FK       <--- category_key PK\n"
    "   brand_name     |    time_key FK            category_name\n"
    "   [SCD Type 1]   |    product_key FK         parent_category\n"
    "                  |    brand_key FK\n"
    "   dim_product    |    category_key FK\n"
    "   product_key PK-+    meal_type, quantity_g\n"
    "   barcode        |    energy_kcal, fat_g, proteins_g, ...\n"
    "   [SCD Type 2]   |\n"
    "                  |    fact_product_market\n"
    "   dim_country    +--> product/brand/category/country_key FK\n"
    "   country_key PK      nutriscore_grade/score, nova_group\n"
    "   [SCD Type 3]        energy/fat/sugar/salt/protein per 100g"
)
add_textbox(slide, 0.8, 1.3, 11, 5.5, schema, 13, False, ACCENT_GREEN, font_name="Courier New")

# E5 - C13 Dimensions & Facts
slide = add_content_slide("E5 | C13 - Dimensions & Fact Tables", BLOCK3_COLOR)
add_table_to_slide(
    slide,
    0.5,
    1.3,
    5.5,
    3.5,
    ["Dimension", "SCD Type", "Key Fields"],
    [
        ["dim_time", "N/A (pre-populated)", "2020-2030, 4018 rows"],
        ["dim_product", "Type 2 (historical)", "barcode, name, ingredients"],
        ["dim_brand", "Type 1 (overwrite)", "brand_name, parent_company"],
        ["dim_category", "None", "name, parent, level"],
        ["dim_country", "Type 3 (add column)", "name, previous_list"],
        ["dim_user", "Type 2 (anonymized)", "user_hash SHA256"],
        ["dim_nutriscore", "None (reference)", "A-E grades, 5 rows"],
    ],
    BLOCK3_COLOR,
    11,
)
add_table_to_slide(
    slide,
    6.5,
    1.3,
    6,
    2.5,
    ["Fact Table", "Grain", "Measures"],
    [
        ["fact_daily_nutrition", "user+day+product+meal", "quantity, kcal, macros, scores"],
        ["fact_product_market", "product+time snapshot", "scores, nutrition per 100g"],
    ],
    BLOCK3_COLOR,
    11,
)
add_textbox(
    slide,
    0.5,
    5.0,
    12,
    0.8,
    "Approach: Bottom-up (Kimball) - build datamarts for known analyses, extend incrementally.\n"
    "3 Datamart views: dm_user_daily_nutrition, dm_product_market_by_category, dm_brand_quality_ranking",
    14,
    False,
    LIGHT_GRAY,
)

# E5 - C14 DW Implementation
slide = add_content_slide("E5 | C14 - DW Implementation & Access Config", BLOCK3_COLOR)
add_bullet_slide(
    slide,
    0.8,
    1.3,
    11,
    5.5,
    [
        ("Implementation: sql/init/02_schema_warehouse.sql (360 lines)", 0),
        ("Separate dw schema within same PostgreSQL instance", 1),
        ("Auto-created during docker compose up (init scripts)", 1),
        ("Access Configuration:", 0),
        ("nutritionist_role: SELECT on datamart views + dim_time + dim_nutriscore only", 1),
        ("etl_service: full DML (SELECT/INSERT/UPDATE/DELETE) on all dw tables", 1),
        ("admin_role: full access to entire warehouse", 1),
        ("Test Procedure:", 0),
        ("Schema creation: all tables, indexes, views created without errors", 1),
        ("dim_time: 4,018 rows (2020-2030) pre-populated", 1),
        ("SCD Type 2: product update creates new row, closes old with end_date", 1),
        ("Datamart views: return aggregated results for analysts", 1),
    ],
    color=WHITE,
    level_colors=[BLOCK3_COLOR, LIGHT_GRAY],
)

# E5 - C15 ETL DAGs
slide = add_content_slide("E5 | C15 - ETL Pipelines (6 Airflow DAGs)", BLOCK3_COLOR)
add_table_to_slide(
    slide,
    0.5,
    1.3,
    12,
    4.0,
    ["DAG", "Schedule", "Tasks", "Competency"],
    [
        ["etl_extract_off_api", "Daily 02:00", "1 task: REST API extraction", "C8"],
        ["etl_extract_parquet", "Weekly Sun 01:00", "2 tasks: Parquet + analytics", "C8, C9"],
        ["etl_extract_scraping", "Weekly Mon 03:00", "1 task: Web scraping", "C8"],
        ["etl_aggregate_clean", "Daily 04:00", "3 tasks: aggregate, clean, load", "C10"],
        ["etl_load_warehouse", "Daily 05:00", "6 tasks: 4 dims + 2 facts", "C13-C17"],
        ["etl_datalake_ingest", "Daily 06:00", "4 tasks: bronze/silver/gold/catalog", "C18-C21"],
    ],
    BLOCK3_COLOR,
    13,
)
etl_flow = (
    "Warehouse ETL task dependencies:\n\n"
    "  [load_dim_brands]     --+\n"
    "  [load_dim_categories] --+--> [load_fact_product_market]\n"
    "  [load_dim_products]   --+--> [load_fact_daily_nutrition]\n"
    "  [load_dim_users]      --+"
)
add_textbox(slide, 0.8, 5.2, 11, 2.0, etl_flow, 13, False, ACCENT_GREEN, font_name="Courier New")

# E5 - C15 ETL Code
slide = add_content_slide("E5 | C15 - ETL Code: Dimension Loading with SCD", BLOCK3_COLOR)
etl_code = (
    "def load_dim_products(**context):\n"
    '    """SCD Type 2: historical tracking"""\n'
    "    with engine.begin() as conn:\n"
    "        # 1. Insert new products\n"
    '        conn.execute(text("""\n'
    "            INSERT INTO dw.dim_product (...)\n"
    "            SELECT p.barcode, p.product_name, ...\n"
    "            FROM app.products p\n"
    "            LEFT JOIN dw.dim_product dp\n"
    "              ON p.barcode = dp.barcode\n"
    "              AND dp.is_current = TRUE\n"
    "            WHERE dp.product_key IS NULL\n"
    '        """))\n'
    "        # 2. Close changed records\n"
    '        conn.execute(text("""\n'
    "            UPDATE dw.dim_product dp\n"
    "            SET end_date = CURRENT_DATE - 1,\n"
    "                is_current = FALSE\n"
    "            FROM app.products p\n"
    "            WHERE dp.barcode = p.barcode\n"
    "              AND dp.is_current = TRUE\n"
    "              AND p.product_name IS DISTINCT\n"
    "                  FROM dp.product_name\n"
    '        """))'
)
add_textbox(slide, 0.8, 1.3, 11, 5.8, etl_code, 12, False, ACCENT_GREEN, font_name="Courier New")

# E5 - C15 Fact loading
slide = add_content_slide("E5 | C15 - ETL: Fact Table Loading", BLOCK3_COLOR)
add_bullet_slide(
    slide,
    0.8,
    1.3,
    11,
    5.5,
    [
        ("load_fact_daily_nutrition:", 0),
        ("Loads yesterday's meal data (incremental, not full reload)", 1),
        ("Joins through 5 dimensions: user, time, product, category, brand", 1),
        ("User pseudonymization: encode(digest(user_id, 'sha256'), 'hex')", 1),
        ("Deduplication: NOT EXISTS prevents duplicate fact rows", 1),
        ("load_fact_product_market:", 0),
        ("Snapshots all products with their current dimension keys", 1),
        ("LEFT JOINs preserve products without brand/category", 1),
        ("Time key: today's date as YYYYMMDD integer", 1),
        ("Output Formats:", 0),
        ("SQL tables -> Datamart views (dm_user_daily_nutrition, etc.)", 1),
        ("Gold layer: Parquet + CSV exports to MinIO", 1),
    ],
    color=WHITE,
    level_colors=[BLOCK3_COLOR, LIGHT_GRAY],
)

# =====================================================================
#  E6 - DW MAINTENANCE (10 min Q&A, ~8 slides)
# =====================================================================
add_section_divider(
    "DW Maintenance & SCD Variations",
    "C16, C17: Administration, supervision, dimension variations",
    BLOCK3_COLOR,
    "E6 - Case Study (Q&A)",
    "10 minutes Q&A | Block 3",
)

# E6 - C16 Logging & Alerts
slide = add_content_slide("E6 | C16 - Activity Logging & Alert System", BLOCK3_COLOR)
add_bullet_slide(
    slide,
    0.8,
    1.3,
    11,
    5.5,
    [
        ("Activity Logging (extraction_log table):", 0),
        ("source_name, source_type, records_extracted/loaded/rejected", 1),
        ("status: running | completed | failed (with error_message)", 1),
        ("Timestamps: started_at, completed_at", 1),
        ("Alert System:", 0),
        ("Airflow: email_on_failure: True -> admin@nutritrack.local", 1),
        ("2 retries with 5-minute delay before alerting", 1),
        ("Task name, error message, execution timestamp in notification", 1),
        ("Service Level Indicators:", 0),
        ("ETL pipeline success rate: target > 95%", 1),
        ("Data freshness: < 24h old", 1),
        ("DW query response time: < 5 seconds (datamart views)", 1),
        ("Backup success rate: 100%", 1),
    ],
    color=WHITE,
    level_colors=[BLOCK3_COLOR, LIGHT_GRAY],
)

# E6 - C16 Backups
slide = add_content_slide("E6 | C16 - Backup Procedures", BLOCK3_COLOR)
add_bullet_slide(
    slide,
    0.8,
    1.3,
    11,
    3.0,
    [
        ("Full backup: make backup (both app + dw schemas)", 0),
        ("Partial backup: make backup-dw (data warehouse only)", 0),
        ("pg_dump --format=custom --compress=6", 0),
        ("Optional upload to MinIO backups bucket (backups/YYYY/MM/DD/)", 0),
        ("Local cleanup: remove backups older than 7 days", 0),
    ],
    color=LIGHT_GRAY,
    font_size=16,
)
add_bullet_slide(
    slide,
    0.8,
    4.5,
    11,
    3.0,
    [
        ("Scalability Documentation:", 0),
        ("New source: add script + DAG + column mappings + catalog update", 1),
        ("New access: CREATE ROLE + GRANT group role + MinIO policy", 1),
        ("Storage increase: Docker volume resize or external PV", 1),
        ("New datamart: CREATE VIEW on dw schema + GRANT to role", 1),
    ],
    color=WHITE,
    level_colors=[BLOCK3_COLOR, LIGHT_GRAY],
)

# E6 - C17 SCD Type 1
slide = add_content_slide("E6 | C17 - SCD Type 1: Overwrite (Brand Corrections)", BLOCK3_COLOR)
scd1_code = (
    "-- Use case: 'Dannon' -> 'Danone' (typo correction)\n\n"
    "CREATE OR REPLACE FUNCTION\n"
    "  dw.scd_type1_update_brand(\n"
    "    p_brand_key INTEGER,\n"
    "    p_new_brand_name VARCHAR(255)\n"
    ") RETURNS VOID AS $$\n"
    "BEGIN\n"
    "    UPDATE dw.dim_brand\n"
    "    SET brand_name = p_new_brand_name,\n"
    "        last_updated = CURRENT_TIMESTAMP\n"
    "    WHERE brand_key = p_brand_key;\n"
    "END;\n"
    "$$ LANGUAGE plpgsql;\n\n"
    "-- Impact: No history preserved.\n"
    "-- Suitable for corrections where old value was wrong."
)
add_textbox(slide, 0.8, 1.3, 11, 5.5, scd1_code, 14, False, ACCENT_GREEN, font_name="Courier New")

# E6 - C17 SCD Type 2
slide = add_content_slide("E6 | C17 - SCD Type 2: Historical (Product Reformulations)", BLOCK3_COLOR)
scd2_code = (
    "-- Use case: Product reformulated, Nutri-Score D -> B\n\n"
    "CREATE OR REPLACE FUNCTION\n"
    "  dw.scd_type2_update_product(\n"
    "    p_barcode VARCHAR(20), ...\n"
    ") RETURNS VOID AS $$\n"
    "BEGIN\n"
    "    -- 1. Close current record\n"
    "    UPDATE dw.dim_product\n"
    "    SET end_date = CURRENT_DATE - 1,\n"
    "        is_current = FALSE\n"
    "    WHERE barcode = p_barcode\n"
    "      AND is_current = TRUE;\n\n"
    "    -- 2. Insert new current record\n"
    "    INSERT INTO dw.dim_product (\n"
    "        barcode, product_name, ...,\n"
    "        effective_date, is_current\n"
    "    ) VALUES (\n"
    "        p_barcode, p_product_name, ...,\n"
    "        CURRENT_DATE, TRUE);\n"
    "END;\n"
    "$$ LANGUAGE plpgsql;"
)
add_textbox(slide, 0.8, 1.3, 11, 5.8, scd2_code, 13, False, ACCENT_GREEN, font_name="Courier New")

# E6 - C17 SCD Type 3
slide = add_content_slide("E6 | C17 - SCD Type 3: Add Column (Country Expansion)", BLOCK3_COLOR)
scd3_code = (
    "-- Use case: 'France' -> 'France, Belgium, Luxembourg'\n\n"
    "CREATE OR REPLACE FUNCTION\n"
    "  dw.scd_type3_update_country(\n"
    "    p_country_key INTEGER,\n"
    "    p_new_country_name VARCHAR(100)\n"
    ") RETURNS VOID AS $$\n"
    "BEGIN\n"
    "    UPDATE dw.dim_country\n"
    "    SET previous_country_list = country_name,\n"
    "        country_name = p_new_country_name\n"
    "    WHERE country_key = p_country_key;\n"
    "END;\n"
    "$$ LANGUAGE plpgsql;\n\n"
    "-- dim_country structure:\n"
    "--   country_name VARCHAR(100)\n"
    "--   previous_country_list VARCHAR(500)  <-- SCD Type 3"
)
add_textbox(slide, 0.8, 1.3, 11, 5.5, scd3_code, 14, False, ACCENT_GREEN, font_name="Courier New")

# E6 - C17 History View
slide = add_content_slide("E6 | C17 - Product History View & Integration", BLOCK3_COLOR)
add_bullet_slide(
    slide,
    0.8,
    1.3,
    11,
    5.5,
    [
        ("History View (v_product_history):", 0),
        ("Shows all versions of products with effective_date and version_status", 1),
        ("SELECT * FROM dw.v_product_history WHERE barcode = '3017620422003'", 1),
        ("Automated Change Detection:", 0),
        ("scd_type2_check_and_update_products() scans app.products vs dw.dim_product", 1),
        ("Uses IS DISTINCT FROM (handles NULLs correctly, unlike <> operator)", 1),
        ("Returns table of (barcode, change_type: 'NEW' or 'UPDATED')", 1),
        ("ETL Integration:", 0),
        ("load_dim_brands: SCD Type 1 (overwrite parent_company changes)", 1),
        ("load_dim_products: SCD Type 2 (close old, insert new current)", 1),
        ("load_dim_users: SCD Type 2 (track profile changes, SHA256 hash)", 1),
        ("Documentation updated: effective_date, end_date, is_current columns documented", 1),
    ],
    color=WHITE,
    level_colors=[BLOCK3_COLOR, LIGHT_GRAY],
)

# =====================================================================
#  E7 - DATA LAKE (10 min, ~10 slides)
# =====================================================================
add_section_divider(
    "Data Lake Architecture & Governance",
    "C18, C19, C20, C21: Architecture, infrastructure, catalog, governance",
    BLOCK4_COLOR,
    "E7 - Professional Situation",
    "10 minutes | Block 4",
)

# E7 - C18 Architecture
slide = add_content_slide("E7 | C18 - Data Lake Architecture (Medallion)", BLOCK4_COLOR)
lake_arch = (
    "+-----------------------------------------------------+\n"
    "|              NutriTrack Data Lake (MinIO)            |\n"
    "|                                                     |\n"
    "|  +-------------+  +-------------+  +-----------+    |\n"
    "|  |   BRONZE    |  |   SILVER    |  |   GOLD    |    |\n"
    "|  | (Raw Data)  |  | (Cleaned)   |  |(Analytics)|    |\n"
    "|  |             |  |             |  |           |    |\n"
    "|  | api/{ds}/   |  | products/   |  | analytics/|    |\n"
    "|  | parquet/    |  |  cleaned.   |  |  nutri-   |    |\n"
    "|  | scraping/   |  |  parquet    |  |  score_   |    |\n"
    "|  | duckdb/     |  |             |  |  dist.    |    |\n"
    "|  | _manifests/ |  | _reports/   |  |  category |    |\n"
    "|  +-------------+  +-------------+  |  brand_   |    |\n"
    "|  Retention:90d    Retention:1yr     |  ranking  |    |\n"
    "|                                    +-----------+    |\n"
    "|  +-------------+                   Retention:indef  |\n"
    "|  |  BACKUPS    |   _catalog/metadata.json           |\n"
    "|  |  YYYY/MM/DD |   (in all buckets)                 |\n"
    "|  +-------------+                                    |\n"
    "+-----------------------------------------------------+"
)
add_textbox(slide, 0.8, 1.3, 11, 5.8, lake_arch, 12, False, ACCENT_GREEN, font_name="Courier New")

# E7 - C18 Volume/Velocity/Variety
slide = add_content_slide("E7 | C18 - Addressing Volume, Velocity, Variety", BLOCK4_COLOR)
add_table_to_slide(
    slide,
    0.8,
    1.3,
    11.5,
    2.5,
    ["V", "Challenge", "Solution"],
    [
        [
            "Volume",
            "3M+ products in OFF dataset (~5 GB)",
            "DuckDB reads Parquet without memory loading; MinIO scales to PB",
        ],
        ["Variety", "JSON, Parquet, CSV, SQL", "Bronze stores raw formats as-is; Silver normalizes to unified Parquet"],
        ["Velocity", "Daily incremental + weekly bulk", "Airflow DAGs on schedule; date-partitioned bronze snapshots"],
    ],
    BLOCK4_COLOR,
    13,
)
add_textbox(slide, 0.8, 4.2, 11, 0.5, "Catalog Tool Comparison:", 20, True, BLOCK4_COLOR)
add_table_to_slide(
    slide,
    0.8,
    4.8,
    11.5,
    2.0,
    ["Criteria", "Apache Atlas", "DataHub", "Custom JSON (selected)"],
    [
        ["Complexity", "High (Hadoop)", "Medium", "Low"],
        ["Setup time", "Days", "Hours", "Minutes"],
        ["Resource cost", "High (JVM)", "Medium", "Minimal"],
        ["Justification", "--", "--", "Minimal overhead, direct MinIO integration, sufficient for project scale"],
    ],
    BLOCK4_COLOR,
    11,
)

# E7 - C19 Infrastructure
slide = add_content_slide("E7 | C19 - Infrastructure Components", BLOCK4_COLOR)
add_bullet_slide(
    slide,
    0.8,
    1.3,
    11,
    5.5,
    [
        ("Storage System: MinIO (S3-compatible)", 0),
        ("Docker image: minio/minio:latest, ports 9000 (API) + 9001 (Console)", 1),
        ("Buckets: bronze, silver, gold, backups (auto-created by minio-init service)", 1),
        ("Healthcheck: curl http://localhost:9000/minio/health/live", 1),
        ("Batch Tool: Apache Airflow 2.8.1", 0),
        ("CeleryExecutor with Redis broker + PostgreSQL backend", 1),
        ("6 DAGs with dependency chains and error handling", 1),
        ("Web UI at port 8080 for monitoring and manual triggers", 1),
        ("Installation Testing (dev environment):", 0),
        ("docker compose up -d -> 10 services running with healthchecks", 1),
        ("make setup-lake -> initialize MinIO buckets + lifecycle rules", 1),
        ("make pipeline -> full extraction -> clean -> import test", 1),
        ("make lake-status -> verify object counts and sizes per bucket", 1),
    ],
    color=WHITE,
    level_colors=[BLOCK4_COLOR, LIGHT_GRAY],
)

# E7 - C19 DAG Pipeline
slide = add_content_slide("E7 | C19 - Medallion Pipeline DAG", BLOCK4_COLOR)
dag_code = (
    "# airflow/dags/etl_datalake_ingest.py\n\n"
    "# Task chain: bronze >> silver >> gold >> catalog\n\n"
    "def ingest_to_bronze(**context):\n"
    "    client = _get_minio_client()\n"
    "    ds = context['ds']  # Airflow execution date\n"
    "    # Upload raw files with date partitioning\n"
    "    for f in api_dir.glob('*.json'):\n"
    "        client.fput_object('bronze',\n"
    "            f'api/{ds}/{f.name}', str(f))\n"
    "    # Write manifest\n"
    "    manifest = {'date': ds, 'files': count}\n"
    "    client.put_object('bronze',\n"
    "        f'_manifests/{ds}.json', ...)\n\n"
    "def transform_to_silver(**context):\n"
    "    # Upload cleaned products to silver\n\n"
    "def publish_to_gold(**context):\n"
    "    # Export 3 analytical datasets:\n"
    "    # nutriscore_distribution, category_stats,\n"
    "    # brand_ranking -> Parquet + CSV\n\n"
    "bronze >> silver >> gold >> catalog"
)
add_textbox(slide, 0.8, 1.3, 11, 5.8, dag_code, 12, False, ACCENT_GREEN, font_name="Courier New")

# E7 - C20 Catalog
slide = add_content_slide("E7 | C20 - Data Catalog Management", BLOCK4_COLOR)
add_bullet_slide(
    slide,
    0.8,
    1.3,
    11,
    5.5,
    [
        ("Catalog stored as _catalog/metadata.json in each bucket", 0),
        ("Updated daily by etl_datalake_ingest DAG (update_catalog_metadata task)", 0),
        ("Per-dataset metadata:", 0),
        ("description, format, update_frequency, source, schema, owner, quality", 1),
        ("lineage tracking (e.g., silver/products <- bronze/api + bronze/parquet)", 1),
        ("Feed Method Selection:", 0),
        ("OFF API: daily pull (API rate limits allow daily)", 1),
        ("OFF Parquet: weekly pull (bulk export updated weekly)", 1),
        ("ANSES/EFSA scraping: weekly (reference values change infrequently)", 1),
        ("Lifecycle & Deletion:", 0),
        ("Bronze api/ and scraping/: auto-expire at 90 days (lifecycle rules)", 1),
        ("Backups daily/: auto-expire at 30 days", 1),
        ("Silver: 1 year retention | Gold: indefinite", 1),
        ("Monitoring: make lake-status -> object count + total size per bucket", 0),
    ],
    color=WHITE,
    level_colors=[BLOCK4_COLOR, LIGHT_GRAY],
)

# E7 - C21 Governance
slide = add_content_slide("E7 | C21 - Data Governance & Access Control", BLOCK4_COLOR)
add_table_to_slide(
    slide,
    0.5,
    1.3,
    12,
    3.5,
    ["Group", "PostgreSQL Role", "MinIO Policy", "API Role"],
    [
        ["End Users", "app_readonly", "gold: readonly", "user"],
        ["Nutritionists", "nutritionist_role", "silver+gold: readonly", "nutritionist"],
        ["Administrators", "admin_role", "all: readwrite", "admin"],
        ["ETL Service", "etl_service", "bronze/silver/gold: readwrite", "N/A (service)"],
    ],
    BLOCK4_COLOR,
    13,
)
add_bullet_slide(
    slide,
    0.8,
    5.2,
    11,
    2.0,
    [
        ("Rights applied to groups (not individuals) when possible", 0),
        ("Access limited to necessary resources per group", 0),
        ("No personal data in data lake (user data stays in PostgreSQL only)", 0),
        ("RGPD: consent tracking, SHA256 pseudonymization, automated cleanup, data registry", 0),
    ],
    color=LIGHT_GRAY,
    font_size=14,
)

# E7 - C21 Group policies code
slide = add_content_slide("E7 | C21 - Group Policies Implementation", BLOCK4_COLOR)
gov_code = (
    "# scripts/setup_minio.py\n\n"
    "GROUP_POLICIES = {\n"
    '    "app_users": {\n'
    '        "description": "End-users: read-only gold",\n'
    '        "buckets": {"gold": "readonly"},\n'
    "    },\n"
    '    "nutritionists": {\n'
    '        "description": "Read silver + gold",\n'
    '        "buckets": {"silver": "readonly",\n'
    '                    "gold": "readonly"},\n'
    "    },\n"
    '    "admins": {\n'
    '        "description": "Full access all buckets",\n'
    '        "buckets": {"bronze": "readwrite",\n'
    '                    "silver": "readwrite",\n'
    '                    "gold": "readwrite",\n'
    '                    "backups": "readwrite"},\n'
    "    },\n"
    '    "etl_service": {\n'
    '        "description": "ETL write access",\n'
    '        "buckets": {"bronze": "readwrite",\n'
    '                    "silver": "readwrite",\n'
    '                    "gold": "readwrite"},\n'
    "    },\n"
    "}"
)
add_textbox(slide, 0.8, 1.3, 11, 5.8, gov_code, 12, False, ACCENT_GREEN, font_name="Courier New")

# =====================================================================
#  COMPETENCY MATRIX
# =====================================================================
slide = add_content_slide("Competency Verification Matrix (C1-C21)", ACCENT_BLUE)
add_table_to_slide(
    slide,
    0.2,
    1.2,
    12.8,
    5.5,
    ["C#", "Competency", "Eval", "Key Artifact"],
    [
        ["C1", "Need analysis & feasibility", "E1,E2", "Interview grids, synthesis note"],
        ["C2", "Data topography & mapping", "E2", "4-part topography, flux matrix"],
        ["C3", "Technical framework design", "E2", "Architecture study, Docker Compose"],
        ["C4", "Technical monitoring", "E2", "Airflow UI, extraction_log"],
        ["C5", "Project planning", "E3", "Roadmap, calendar, effort weighting"],
        ["C6", "Project supervision", "E2,E3", "Makefile, Airflow UI, indicators"],
        ["C7", "Communication strategy", "E3", "Comms plan, Swagger, Streamlit"],
        ["C8", "Data extraction scripts", "E4", "5 scripts (API, Parquet, scraping, DB, DuckDB)"],
        ["C9", "SQL queries", "E4", "7 queries with EXPLAIN optimization notes"],
        ["C10", "Aggregation & cleaning", "E4", "aggregate_clean.py + cleaning report"],
        ["C11", "RGPD-compliant database", "E4", "MCD/MLD/MPD, RGPD registry, import"],
        ["C12", "REST API", "E4", "FastAPI, JWT, 10+ endpoints, OpenAPI"],
        ["C13", "Star schema modeling", "E5", "2 facts, 7 dimensions, bottom-up"],
        ["C14", "DW implementation", "E5", "dw schema, access config, tests"],
        ["C15", "ETL pipelines", "E5", "6 Airflow DAGs"],
        ["C16", "DW maintenance", "E6", "Logging, alerts, backups, docs"],
        ["C17", "SCD variations", "E6", "Type 1/2/3 procedures + history view"],
        ["C18", "Data lake architecture", "E7", "Medallion, V/V/V, catalog comparison"],
        ["C19", "Infrastructure components", "E7", "MinIO + Airflow in Docker Compose"],
        ["C20", "Data catalog", "E7", "metadata.json, lifecycle, monitoring"],
        ["C21", "Data governance", "E7", "Group policies, RBAC, RGPD compliance"],
    ],
    ACCENT_BLUE,
    10,
)

# =====================================================================
#  THANK YOU / Q&A SLIDE
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_textbox(slide, 1, 2.0, 11, 1.5, "Thank You", 54, True, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 1, 3.5, 11, 0.8, "Questions?", 36, False, ACCENT_BLUE, PP_ALIGN.CENTER)
add_textbox(
    slide, 1, 5.0, 11, 0.5, "NutriTrack | RNCP37638 Data Engineer Certification", 18, False, LIGHT_GRAY, PP_ALIGN.CENTER
)
add_textbox(
    slide,
    1,
    5.6,
    11,
    0.5,
    "Git Repository  |  docs/rapport_final.md  |  docker compose up -d",
    16,
    False,
    MEDIUM_GRAY,
    PP_ALIGN.CENTER,
)

# ── Save ─────────────────────────────────────────────────────────────
output_path = "/Users/reegauta/Documents/Simplon/nutritrack/docs/NutriTrack_Defense_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
