#!/usr/bin/env python3
"""
Generate NutriTrack Defense Presentation (PPTX) from defense_slides.tex content.

Produces a professional 16:9 PPTX with green theme matching NutriTrack branding.
Colors: #2E7D32 primary green, #FF8F00 accent orange, #1565C0 blue.

Output: NutriTrack_Defense_Visual.pptx
"""

from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Theme colors
# ---------------------------------------------------------------------------
GREEN = RGBColor(0x2E, 0x7D, 0x32)
LIGHT_GREEN = RGBColor(0x66, 0xBB, 0x6A)
DARK_GREEN = RGBColor(0x1B, 0x5E, 0x20)
ACCENT = RGBColor(0xFF, 0x8F, 0x00)
BLUE = RGBColor(0x15, 0x65, 0xC0)
RED = RGBColor(0xC6, 0x28, 0x28)
PURPLE = RGBColor(0x6A, 0x1B, 0x9A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
BG_WHITE = RGBColor(0xFA, 0xFA, 0xFA)
BRONZE = RGBColor(0xCD, 0x7F, 0x32)
SILVER = RGBColor(0xA0, 0xA0, 0xA0)
GOLD = RGBColor(0xDA, 0xA5, 0x20)

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Fonts
TITLE_FONT = "Calibri"
BODY_FONT = "Calibri"

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------


def add_slide_number(slide, num):
    """Add slide number in bottom-right corner."""
    txBox = slide.shapes.add_textbox(SLIDE_WIDTH - Inches(0.8), SLIDE_HEIGHT - Inches(0.45), Inches(0.6), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT


def set_slide_bg(slide, color):
    """Set slide background to a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title_text, subtitle_text=None, competency=None):
    """Add a green title bar at the top of the slide."""
    # Green rectangle
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), SLIDE_WIDTH - Inches(1), Inches(0.55))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = TITLE_FONT

    if competency:
        run2 = p.add_run()
        run2.text = f"  {competency}"
        run2.font.size = Pt(16)
        run2.font.bold = True
        run2.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)  # light blue
        run2.font.name = TITLE_FONT

    # Accent line under title bar
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), SLIDE_WIDTH, Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def add_bullet_list(
    slide, items, left, top, width, height, font_size=14, color=DARK_GRAY, bold_prefix=False, bullet_color=GREEN
):
    """Add a bullet list to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        p.space_before = Pt(2)
        p.level = 0

        # Add bullet character
        run_bullet = p.add_run()
        run_bullet.text = "\u2022  "
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = BODY_FONT

        if bold_prefix and ": " in item:
            prefix, rest = item.split(": ", 1)
            run_b = p.add_run()
            run_b.text = prefix + ": "
            run_b.font.size = Pt(font_size)
            run_b.font.bold = True
            run_b.font.color.rgb = color
            run_b.font.name = BODY_FONT
            run_t = p.add_run()
            run_t.text = rest
            run_t.font.size = Pt(font_size)
            run_t.font.color.rgb = color
            run_t.font.name = BODY_FONT
        else:
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = BODY_FONT
    return txBox


def add_textbox(
    slide,
    text,
    left,
    top,
    width,
    height,
    font_size=14,
    color=DARK_GRAY,
    bold=False,
    alignment=PP_ALIGN.LEFT,
    font_name=None,
):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name or BODY_FONT
    return txBox


def add_rounded_box(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title, items, fill_color, border_color, title_color=None, icon_char=None):
    """Add a card with title and bullet items."""
    box = add_rounded_box(slide, left, top, width, height, fill_color, border_color)
    # Title
    y_offset = top + Inches(0.15)
    title_text = (icon_char + "  " if icon_char else "") + title
    add_textbox(
        slide,
        title_text,
        left + Inches(0.15),
        y_offset,
        width - Inches(0.3),
        Inches(0.35),
        font_size=14,
        color=title_color or border_color,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    # Items
    if items:
        add_bullet_list(
            slide,
            items,
            left + Inches(0.15),
            y_offset + Inches(0.4),
            width - Inches(0.3),
            height - Inches(0.6),
            font_size=11,
            color=DARK_GRAY,
            bullet_color=border_color,
        )
    return box


def add_demo_callout(slide, text, top):
    """Add a LIVE DEMO callout box."""
    left = Inches(2.5)
    width = SLIDE_WIDTH - Inches(5)
    add_rounded_box(slide, left, top, width, Inches(0.55), RGBColor(0xFF, 0xF3, 0xE0), ACCENT)
    add_textbox(
        slide,
        f"LIVE DEMO: {text}",
        left + Inches(0.2),
        top + Inches(0.08),
        width - Inches(0.4),
        Inches(0.4),
        font_size=13,
        color=ACCENT,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )


def add_why_callout(slide, text, top):
    """Add a 'Why?' callout box."""
    left = Inches(0.5)
    width = SLIDE_WIDTH - Inches(1)
    add_rounded_box(slide, left, top, width, Inches(0.6), RGBColor(0xE3, 0xF2, 0xFD), BLUE)
    add_textbox(
        slide,
        f"Why? {text}",
        left + Inches(0.2),
        top + Inches(0.05),
        width - Inches(0.4),
        Inches(0.5),
        font_size=11,
        color=BLUE,
        bold=False,
    )


def add_table(slide, left, top, width, height, headers, rows):
    """Add a properly formatted table."""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    # Distribute column widths
    col_width = int(width / num_cols)
    for i in range(num_cols):
        table.columns[i].width = col_width

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = BODY_FONT
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN

    # Data rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.color.rgb = DARK_GRAY
                paragraph.font.name = BODY_FONT
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    return table_shape


def add_flow_arrow(slide, x1, y1, x2, y2, color=GRAY):
    """Add a connector arrow between two points."""
    connector = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, min(x1, x2), min(y1, y2), abs(x2 - x1) or Inches(0.4), abs(y2 - y1) or Inches(0.15)
    )
    connector.fill.solid()
    connector.fill.fore_color.rgb = color
    connector.line.fill.background()
    return connector


def add_section_slide(slide, title, subtitle=None, bg_color=GREEN):
    """Create a section divider slide."""
    set_slide_bg(slide, bg_color)
    add_textbox(
        slide,
        title,
        Inches(1),
        Inches(2.5),
        SLIDE_WIDTH - Inches(2),
        Inches(1.2),
        font_size=36,
        color=WHITE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    if subtitle:
        add_textbox(
            slide,
            subtitle,
            Inches(1),
            Inches(3.7),
            SLIDE_WIDTH - Inches(2),
            Inches(0.8),
            font_size=20,
            color=RGBColor(0xC8, 0xE6, 0xC9),
            bold=False,
            alignment=PP_ALIGN.CENTER,
        )


def lighten(color, factor=0.85):
    """Return a lighter version of an RGB color."""
    r = min(255, int(color[0] + (255 - color[0]) * factor))
    g = min(255, int(color[1] + (255 - color[1]) * factor))
    b = min(255, int(color[2] + (255 - color[2]) * factor))
    return RGBColor(r, g, b)


slide_num = 0

# ===================================================================
# SLIDE 1: TITLE
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, GREEN)

# NT logo circle
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.666), Inches(1.0), Inches(2.0), Inches(2.0))
circle.fill.solid()
circle.fill.fore_color.rgb = WHITE
circle.line.fill.background()
tf = circle.text_frame
tf.word_wrap = False
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "NT"
run.font.size = Pt(48)
run.font.bold = True
run.font.color.rgb = GREEN
run.font.name = TITLE_FONT

# Title
add_textbox(
    slide,
    "NutriTrack",
    Inches(1),
    Inches(3.2),
    SLIDE_WIDTH - Inches(2),
    Inches(1),
    font_size=54,
    color=WHITE,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)
add_textbox(
    slide,
    "Nutritional Data Engineering Platform",
    Inches(1),
    Inches(4.2),
    SLIDE_WIDTH - Inches(2),
    Inches(0.6),
    font_size=24,
    color=RGBColor(0xC8, 0xE6, 0xC9),
    bold=False,
    alignment=PP_ALIGN.CENTER,
)
add_textbox(
    slide,
    "Reetika Gautam",
    Inches(1),
    Inches(5.2),
    SLIDE_WIDTH - Inches(2),
    Inches(0.5),
    font_size=20,
    color=WHITE,
    bold=False,
    alignment=PP_ALIGN.CENTER,
)
add_textbox(
    slide,
    "RNCP37638 -- Expert in Massive Data Infrastructures\nLevel 7  |  April 2026  |  Simplon.co",
    Inches(1),
    Inches(5.8),
    SLIDE_WIDTH - Inches(2),
    Inches(0.8),
    font_size=14,
    color=RGBColor(0xA5, 0xD6, 0xA7),
    bold=False,
    alignment=PP_ALIGN.CENTER,
)
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 2: AGENDA
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Defense Agenda -- 90 Minutes")

# Four blocks as boxes with arrows
blocks = [
    ("E1-E5", "Project &\nImplementation", "60 min", "C1-C15", GREEN),
    ("E6", "DW Maintenance\nQ&A", "10 min", "C16-C17", BLUE),
    ("E7", "Data Lake", "10 min", "C18-C21", PURPLE),
    ("Jury", "Questions", "10 min", "", ACCENT),
]

start_x = Inches(0.8)
box_w = Inches(2.6)
box_h = Inches(1.6)
gap = Inches(0.5)
arrow_w = Inches(0.4)
y_pos = Inches(2.5)

for i, (title, desc, time, comps, color) in enumerate(blocks):
    x = start_x + i * (box_w + gap + arrow_w)
    box = add_rounded_box(slide, x, y_pos, box_w, box_h, WHITE, color)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = TITLE_FONT
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = desc
    r2.font.size = Pt(11)
    r2.font.color.rgb = GRAY
    r2.font.name = BODY_FONT

    # Time label below
    add_textbox(
        slide,
        time,
        x,
        y_pos + box_h + Inches(0.1),
        box_w,
        Inches(0.3),
        font_size=12,
        color=GRAY,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    # Competencies above
    if comps:
        add_textbox(
            slide,
            comps,
            x,
            y_pos - Inches(0.35),
            box_w,
            Inches(0.3),
            font_size=10,
            color=color,
            bold=False,
            alignment=PP_ALIGN.CENTER,
        )

    # Arrow between blocks
    if i < len(blocks) - 1:
        arrow_x = x + box_w + Inches(0.05)
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, arrow_x, y_pos + box_h / 2 - Inches(0.12), arrow_w - Inches(0.1), Inches(0.24)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = LIGHT_GRAY
        arrow.line.fill.background()

# Bottom bar
add_textbox(
    slide,
    "Live demos throughout   |   Report handed to jury   |   Code on GitHub",
    Inches(1),
    Inches(5.5),
    SLIDE_WIDTH - Inches(2),
    Inches(0.4),
    font_size=13,
    color=GRAY,
    bold=False,
    alignment=PP_ALIGN.CENTER,
)
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 3: THE PROBLEM
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "The Problem")

problems = [
    ("Scattered", "Data in 5+ sources\nNo single view", RED),
    ("Inconsistent", "Different formats\nMissing fields, dupes", RED),
    ("No Compliance", "No RGPD registry\nNo access control", RED),
]

card_w = Inches(3.5)
card_h = Inches(2.2)
start_x = Inches(0.9)
gap = Inches(0.7)
y_pos = Inches(2.0)

for i, (title, desc, color) in enumerate(problems):
    x = start_x + i * (card_w + gap)
    light_red = RGBColor(0xFC, 0xE4, 0xE4)
    box = add_rounded_box(slide, x, y_pos, card_w, card_h, light_red, RED)
    add_textbox(
        slide,
        title,
        x,
        y_pos + Inches(0.3),
        card_w,
        Inches(0.4),
        font_size=20,
        color=RED,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        desc,
        x,
        y_pos + Inches(0.8),
        card_w,
        Inches(1.0),
        font_size=12,
        color=DARK_GRAY,
        bold=False,
        alignment=PP_ALIGN.CENTER,
    )

add_textbox(
    slide,
    "Nutritionists, analysts, and users cannot trust or use the data.",
    Inches(1),
    Inches(5.0),
    SLIDE_WIDTH - Inches(2),
    Inches(0.5),
    font_size=20,
    color=RED,
    bold=False,
    alignment=PP_ALIGN.CENTER,
)
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 4: THE SOLUTION
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "The Solution -- NutriTrack")

solutions = [
    ("5 Extractors", "API, Parquet, Scraping\nDatabase, DuckDB", GREEN),
    ("Unified Storage", "Star-schema DW\nMedallion data lake", GREEN),
    ("Full RGPD", "Data registry\nConsent + cleanup", GREEN),
]

card_w = Inches(3.5)
card_h = Inches(2.0)
start_x = Inches(0.9)
gap = Inches(0.7)
y_pos = Inches(1.8)

for i, (title, desc, color) in enumerate(solutions):
    x = start_x + i * (card_w + gap)
    light_green = RGBColor(0xE8, 0xF5, 0xE9)
    box = add_rounded_box(slide, x, y_pos, card_w, card_h, light_green, GREEN)
    add_textbox(
        slide,
        title,
        x,
        y_pos + Inches(0.25),
        card_w,
        Inches(0.4),
        font_size=20,
        color=GREEN,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        desc,
        x,
        y_pos + Inches(0.75),
        card_w,
        Inches(0.8),
        font_size=12,
        color=DARK_GRAY,
        bold=False,
        alignment=PP_ALIGN.CENTER,
    )

# Stats bar
stats = "3M+ products   |   15 services   |   7 DAGs   |   6 dashboards   |   1 command to deploy"
stats_bg = add_rounded_box(
    slide, Inches(1.5), Inches(4.5), SLIDE_WIDTH - Inches(3), Inches(0.55), RGBColor(0xE8, 0xF5, 0xE9), GREEN
)
add_textbox(
    slide,
    stats,
    Inches(1.5),
    Inches(4.55),
    SLIDE_WIDTH - Inches(3),
    Inches(0.45),
    font_size=13,
    color=GREEN,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 5: END-TO-END DATA FLOW
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "End-to-End Data Flow")

# Row 1: Extract -> Raw -> Clean -> App DB
flow_items_r1 = [
    ("Extract\n5 scripts", GREEN),
    ("Raw Files\n/data/raw/", ACCENT),
    ("Clean\naggregate_clean", PURPLE),
    ("App DB\nPostgreSQL", BLUE),
]
bw = Inches(2.2)
bh = Inches(1.1)
sx = Inches(0.5)
gp = Inches(0.6)
y1 = Inches(1.6)

for i, (label, color) in enumerate(flow_items_r1):
    x = sx + i * (bw + gp)
    light = lighten((color[0], color[1], color[2]), 0.85)
    box = add_rounded_box(slide, x, y1, bw, bh, light, color)
    add_textbox(
        slide,
        label,
        x,
        y1 + Inches(0.15),
        bw,
        bh - Inches(0.3),
        font_size=12,
        color=color,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    if i < len(flow_items_r1) - 1:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, x + bw + Inches(0.05), y1 + bh / 2 - Inches(0.1), gp - Inches(0.1), Inches(0.2)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = LIGHT_GRAY
        arrow.line.fill.background()

# Row 2: API, Warehouse, Lake
flow_items_r2 = [
    ("FastAPI\nREST", ACCENT, Inches(0.5)),
    ("Warehouse\nStar schema", BLUE, Inches(3.6)),
    ("Data Lake\nMedallion", BRONZE, Inches(9.5)),
]
y2 = Inches(3.4)

for label, color, x in flow_items_r2:
    light = lighten((color[0], color[1], color[2]), 0.85)
    box = add_rounded_box(slide, x, y2, Inches(2.2), bh, light, color)
    add_textbox(
        slide,
        label,
        x,
        y2 + Inches(0.15),
        Inches(2.2),
        bh - Inches(0.3),
        font_size=12,
        color=color,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

# Row 3: Streamlit, Superset, Grafana
flow_items_r3 = [
    ("Streamlit\nFrontend", ACCENT, Inches(0.5)),
    ("Superset\nBI", ACCENT, Inches(3.6)),
    ("Grafana\nMonitoring", RED, Inches(6.7)),
]
y3 = Inches(5.2)

for label, color, x in flow_items_r3:
    light = lighten((color[0], color[1], color[2]), 0.85)
    box = add_rounded_box(slide, x, y3, Inches(2.2), bh, light, color)
    add_textbox(
        slide,
        label,
        x,
        y3 + Inches(0.15),
        Inches(2.2),
        bh - Inches(0.3),
        font_size=12,
        color=color,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

# Users row
users = [
    ("End users", Inches(0.5)),
    ("Analysts", Inches(3.6)),
    ("Ops team", Inches(6.7)),
]
for label, x in users:
    add_textbox(
        slide,
        label,
        x,
        y3 + bh + Inches(0.05),
        Inches(2.2),
        Inches(0.25),
        font_size=10,
        color=GRAY,
        bold=False,
        alignment=PP_ALIGN.CENTER,
    )

# Down arrows from App DB to row 2
for target_x in [Inches(0.5), Inches(3.6), Inches(9.5)]:
    darrow = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        sx + 3 * (bw + gp) + bw / 2 - Inches(0.1),
        y1 + bh + Inches(0.05),
        Inches(0.2),
        Inches(0.3),
    )
    darrow.fill.solid()
    darrow.fill.fore_color.rgb = LIGHT_GRAY
    darrow.line.fill.background()
    break  # Only one central arrow

add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 6: DOCKER - ONE COMMAND, 15 SERVICES
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "One Command, 15 Services")

add_textbox(
    slide,
    "docker compose up -d",
    Inches(3),
    Inches(1.4),
    Inches(7),
    Inches(0.5),
    font_size=24,
    color=GREEN,
    bold=True,
    alignment=PP_ALIGN.CENTER,
    font_name="Courier New",
)

categories = [
    ("CORE (4)", BLUE, ["PostgreSQL", "Redis", "MinIO", "MinIO-init"]),
    ("AIRFLOW (4)", PURPLE, ["Webserver", "Scheduler", "Worker", "Init"]),
    ("APPS (3)", ACCENT, ["FastAPI", "Superset", "Streamlit"]),
    ("MONITOR (4)", RED, ["Prometheus", "Grafana", "StatsD-exp", "MailHog"]),
]

cat_x = Inches(0.4)
cat_w = Inches(2.9)
cat_h = Inches(2.3)
cat_gap = Inches(0.3)
y_cat = Inches(2.3)

for ci, (cat_title, color, services) in enumerate(categories):
    x = cat_x + ci * (cat_w + cat_gap)
    # Category header
    header = add_rounded_box(slide, x, y_cat, cat_w, Inches(0.4), color)
    add_textbox(
        slide,
        cat_title,
        x,
        y_cat + Inches(0.03),
        cat_w,
        Inches(0.35),
        font_size=12,
        color=WHITE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    # Service boxes
    light = lighten((color[0], color[1], color[2]), 0.88)
    for si, svc in enumerate(services):
        row = si // 2
        col = si % 2
        sx_inner = x + col * (cat_w / 2)
        sy = y_cat + Inches(0.5) + row * Inches(0.55)
        sw = cat_w / 2 - Inches(0.05)
        sbox = add_rounded_box(slide, sx_inner, sy, sw, Inches(0.45), light, color)
        add_textbox(
            slide,
            svc,
            sx_inner,
            sy + Inches(0.05),
            sw,
            Inches(0.35),
            font_size=9,
            color=color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

# Why callout
add_why_callout(
    slide,
    'Docker Compose = reproducible environment. Any machine with Docker can run the full platform. No "works on my machine" issues.',
    Inches(5.2),
)
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 7: TECH STACK
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Tech Stack -- What & Why")

layers = [
    (
        "EXTRACT",
        GREEN,
        ["OFF API (REST + Parquet)", "BeautifulSoup (Scraping)", "DuckDB (Big data SQL)"],
        "3 source types per C8 req.",
    ),
    ("ORCHESTRATE", PURPLE, ["Apache Airflow 2.8 -- 7 DAGs, Celery, Scheduler"], "vs Prefect: more mature, better UI"),
    (
        "STORE",
        BLUE,
        ["PostgreSQL 16 (OLTP + DW)", "MinIO (S3 Data Lake)", "Redis 7 (Cache + Broker)"],
        "PG: ACID + RGPD. MinIO: free S3",
    ),
    (
        "SERVE",
        ACCENT,
        ["FastAPI (Async + JWT)", "Superset 6 (BI Dashboards)", "Streamlit (User Frontend)"],
        "3 audiences: dev, analyst, user",
    ),
    ("MONITOR", RED, ["Prometheus + Grafana (6 dashboards)", "MailHog (SMTP alerts)"], "C16: alerting + SLA tracking"),
]

lbl_x = Inches(0.3)
lbl_w = Inches(1.6)
items_x = Inches(2.0)
why_x = Inches(9.5)
why_w = Inches(3.2)
row_h = Inches(0.9)
start_y = Inches(1.5)

for i, (label, color, tools, why) in enumerate(layers):
    y = start_y + i * row_h

    # Label box
    lbl_box = add_rounded_box(slide, lbl_x, y + Inches(0.1), lbl_w, Inches(0.4), color)
    add_textbox(
        slide,
        label,
        lbl_x,
        y + Inches(0.13),
        lbl_w,
        Inches(0.35),
        font_size=11,
        color=WHITE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # Tools
    tool_text = "   |   ".join(tools)
    add_textbox(
        slide, tool_text, items_x, y + Inches(0.1), Inches(7.2), Inches(0.4), font_size=11, color=DARK_GRAY, bold=False
    )

    # Why
    add_textbox(slide, why, why_x, y + Inches(0.1), why_w, Inches(0.5), font_size=9, color=GRAY, bold=False)

add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 8: KEY TECHNICAL DECISIONS (C3)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Key Technical Decisions", competency="C3")

decisions = [
    ("PostgreSQL 16", "ACID + RGPD delete\nOLTP + OLAP", "MongoDB", "No ACID, hard RGPD"),
    ("Airflow 2.8", "DAG UI, mature\nCelery for scale", "Prefect", "Newer, less proven"),
    ("MinIO (S3)", "Free, self-hosted\nS3-compatible", "AWS S3", "Cloud cost, lock-in"),
    ("Superset 6", "Open-source BI\nRBAC built-in", "Metabase", "Less enterprise"),
]

card_w = Inches(2.8)
card_h_chosen = Inches(1.8)
card_h_rejected = Inches(1.2)
sx = Inches(0.5)
gp = Inches(0.3)
y_chosen = Inches(1.6)
y_vs = y_chosen + card_h_chosen + Inches(0.1)
y_rejected = y_vs + Inches(0.35)

for i, (chosen, chosen_desc, rejected, rejected_desc) in enumerate(decisions):
    x = sx + i * (card_w + gp)
    # "Selected" label
    add_textbox(
        slide,
        "Selected",
        x,
        y_chosen - Inches(0.25),
        card_w,
        Inches(0.25),
        font_size=9,
        color=GREEN,
        bold=False,
        alignment=PP_ALIGN.CENTER,
    )
    # Chosen box
    light_green = RGBColor(0xE8, 0xF5, 0xE9)
    box = add_rounded_box(slide, x, y_chosen, card_w, card_h_chosen, light_green, GREEN)
    add_textbox(
        slide,
        chosen,
        x,
        y_chosen + Inches(0.15),
        card_w,
        Inches(0.35),
        font_size=14,
        color=GREEN,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        chosen_desc,
        x,
        y_chosen + Inches(0.55),
        card_w,
        Inches(0.9),
        font_size=10,
        color=DARK_GRAY,
        bold=False,
        alignment=PP_ALIGN.CENTER,
    )

    # VS label
    vs_box = add_rounded_box(
        slide, x + card_w / 2 - Inches(0.3), y_vs, Inches(0.6), Inches(0.3), RGBColor(0xFF, 0xF8, 0xE1), ACCENT
    )
    add_textbox(
        slide,
        "vs",
        x + card_w / 2 - Inches(0.3),
        y_vs + Inches(0.02),
        Inches(0.6),
        Inches(0.25),
        font_size=10,
        color=ACCENT,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # Rejected box
    light_gray_bg = RGBColor(0xF5, 0xF5, 0xF5)
    box2 = add_rounded_box(slide, x, y_rejected, card_w, card_h_rejected, light_gray_bg, LIGHT_GRAY)
    add_textbox(
        slide,
        rejected,
        x,
        y_rejected + Inches(0.1),
        card_w,
        Inches(0.3),
        font_size=12,
        color=GRAY,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        rejected_desc,
        x,
        y_rejected + Inches(0.45),
        card_w,
        Inches(0.5),
        font_size=9,
        color=GRAY,
        bold=False,
        alignment=PP_ALIGN.CENTER,
    )

add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 9: THREE SCHEMAS / ZONES
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "PostgreSQL -- Three Zones, One Database")

zones = [
    (
        "App Schema",
        "The Live Application\n\nproducts, users, meals\nbrands, categories\n\nUsed by: FastAPI + Streamlit",
        ACCENT,
    ),
    ("DW Schema", "The Warehouse\n\n7 dims + 2 facts\n6 datamart views\n\nUsed by: Superset", BLUE),
    ("MinIO", "The Data Lake\n\nbronze / silver / gold\nParquet files only\n\nUsed by: Data scientists", PURPLE),
]

card_w = Inches(3.5)
card_h = Inches(3.2)
sx = Inches(0.6)
gp = Inches(0.6)
y_pos = Inches(1.8)

for i, (title, desc, color) in enumerate(zones):
    x = sx + i * (card_w + gp)
    light = lighten((color[0], color[1], color[2]), 0.9)
    box = add_rounded_box(slide, x, y_pos, card_w, card_h, light, color)
    add_textbox(
        slide,
        title,
        x,
        y_pos + Inches(0.2),
        card_w,
        Inches(0.4),
        font_size=18,
        color=color,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        desc,
        x + Inches(0.2),
        y_pos + Inches(0.7),
        card_w - Inches(0.4),
        card_h - Inches(1.0),
        font_size=11,
        color=DARK_GRAY,
        bold=False,
        alignment=PP_ALIGN.CENTER,
    )

# ETL arrows between zones
for i in range(2):
    arrow_x = sx + (i + 1) * card_w + i * gp + Inches(0.1)
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, arrow_x, y_pos + card_h / 2 - Inches(0.1), gp - Inches(0.2), Inches(0.2)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = LIGHT_GRAY
    arrow.line.fill.background()
    add_textbox(
        slide,
        "ETL",
        arrow_x,
        y_pos + card_h / 2 - Inches(0.35),
        gp,
        Inches(0.2),
        font_size=9,
        color=GRAY,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 10: DW vs LAKE
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Why Both a Warehouse AND a Lake?")

# DW column
dw_items = [
    "Product data + User data",
    "ACID transactions",
    "Row-level RGPD delete",
    "Consent filtering per query",
    "<100ms dashboard queries",
    "For BI analysts, dashboards",
]
lake_items = [
    "Product data ONLY",
    "No user data -- ever",
    "Parquet files (immutable)",
    "Schema flexibility",
    "Bulk reads (ML training)",
    "For data scientists, ML",
]

col_w = Inches(5.2)
col_h = Inches(3.2)
dw_x = Inches(0.6)
lake_x = Inches(6.5)
y_pos = Inches(1.6)

# DW box
light_blue = RGBColor(0xE3, 0xF2, 0xFD)
add_rounded_box(slide, dw_x, y_pos, col_w, col_h, light_blue, BLUE)
add_textbox(
    slide,
    "Data Warehouse (PostgreSQL)",
    dw_x + Inches(0.2),
    y_pos + Inches(0.15),
    col_w - Inches(0.4),
    Inches(0.35),
    font_size=16,
    color=BLUE,
    bold=True,
)
add_bullet_list(
    slide,
    dw_items,
    dw_x + Inches(0.3),
    y_pos + Inches(0.6),
    col_w - Inches(0.6),
    col_h - Inches(0.8),
    font_size=12,
    color=DARK_GRAY,
    bullet_color=BLUE,
)

# Lake box
light_purple = RGBColor(0xF3, 0xE5, 0xF5)
add_rounded_box(slide, lake_x, y_pos, col_w, col_h, light_purple, PURPLE)
add_textbox(
    slide,
    "Data Lake (MinIO)",
    lake_x + Inches(0.2),
    y_pos + Inches(0.15),
    col_w - Inches(0.4),
    Inches(0.35),
    font_size=16,
    color=PURPLE,
    bold=True,
)
add_bullet_list(
    slide,
    lake_items,
    lake_x + Inches(0.3),
    y_pos + Inches(0.6),
    col_w - Inches(0.6),
    col_h - Inches(0.8),
    font_size=12,
    color=DARK_GRAY,
    bullet_color=PURPLE,
)

# RGPD boundary
add_textbox(
    slide,
    "RGPD Boundary: User data never enters the lake. Parquet files are immutable -- you cannot delete one user's rows.",
    Inches(1.5),
    Inches(5.3),
    SLIDE_WIDTH - Inches(3),
    Inches(0.5),
    font_size=12,
    color=RED,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)
add_slide_number(slide, slide_num)

# ===================================================================
# SECTION: E1 - Need Analysis (C1)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_slide(slide, "E1 -- Need Analysis", "Competency C1")
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 12: INTERVIEW GRIDS (C1)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Interview Grids", competency="C1")

# Producers card
prod_items = [
    "Business activities",
    "Data types & formats",
    "Volumes & frequency",
    "Quality controls",
    "Metadata available",
    "Access constraints",
    "Known issues",
]
cons_items = [
    "Analysis objectives",
    "Required granularity",
    "Delivery format needs",
    "Frequency needs",
    "RGPD constraints",
    "Tool preferences",
    "Accessibility needs",
]

card_w = Inches(5.2)
card_h = Inches(4.0)
y_pos = Inches(1.8)

# Producers
light_green_bg = RGBColor(0xE8, 0xF5, 0xE9)
add_rounded_box(slide, Inches(0.6), y_pos, card_w, card_h, light_green_bg, GREEN)
# Circle icon
circle_p = slide.shapes.add_shape(
    MSO_SHAPE.OVAL, Inches(0.6) + card_w / 2 - Inches(0.35), y_pos + Inches(0.15), Inches(0.7), Inches(0.7)
)
circle_p.fill.solid()
circle_p.fill.fore_color.rgb = GREEN
circle_p.line.fill.background()
add_textbox(
    slide,
    "Data Producers",
    Inches(0.6),
    y_pos + Inches(0.9),
    card_w,
    Inches(0.35),
    font_size=16,
    color=GREEN,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)
add_bullet_list(
    slide,
    prod_items,
    Inches(1.0),
    y_pos + Inches(1.35),
    card_w - Inches(0.8),
    card_h - Inches(1.5),
    font_size=12,
    color=DARK_GRAY,
    bullet_color=GREEN,
)

# Consumers
light_blue_bg = RGBColor(0xE3, 0xF2, 0xFD)
add_rounded_box(slide, Inches(6.5), y_pos, card_w, card_h, light_blue_bg, BLUE)
circle_c = slide.shapes.add_shape(
    MSO_SHAPE.OVAL, Inches(6.5) + card_w / 2 - Inches(0.35), y_pos + Inches(0.15), Inches(0.7), Inches(0.7)
)
circle_c.fill.solid()
circle_c.fill.fore_color.rgb = BLUE
circle_c.line.fill.background()
add_textbox(
    slide,
    "Data Consumers",
    Inches(6.5),
    y_pos + Inches(0.9),
    card_w,
    Inches(0.35),
    font_size=16,
    color=BLUE,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)
add_bullet_list(
    slide,
    cons_items,
    Inches(6.9),
    y_pos + Inches(1.35),
    card_w - Inches(0.8),
    card_h - Inches(1.5),
    font_size=12,
    color=DARK_GRAY,
    bullet_color=BLUE,
)

add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 13: SMART OBJECTIVES (C1)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "SMART Objectives", competency="C1")

smart = [
    ("S", "Centralize\nfrom 5+ sources"),
    ("M", "50k+ products\n<5s queries"),
    ("A", "Open-source\nDocker-based"),
    ("R", "Fix fragmented\nnutrition data"),
    ("T", "12-week\ntimeline"),
]

sx = Inches(0.5)
bw = Inches(2.2)
bh = Inches(1.6)
gp = Inches(0.35)
y_pos = Inches(2.0)

for i, (letter, desc) in enumerate(smart):
    x = sx + i * (bw + gp)

    # Letter circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x + bw / 2 - Inches(0.3), y_pos - Inches(0.4), Inches(0.6), Inches(0.6)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = GREEN
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = letter
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = TITLE_FONT

    # Description box
    light_green = RGBColor(0xE8, 0xF5, 0xE9)
    box = add_rounded_box(slide, x, y_pos + Inches(0.3), bw, bh, light_green, GREEN)
    add_textbox(
        slide,
        desc,
        x,
        y_pos + Inches(0.5),
        bw,
        bh - Inches(0.3),
        font_size=13,
        color=DARK_GRAY,
        bold=False,
        alignment=PP_ALIGN.CENTER,
    )

    # Arrow
    if i < len(smart) - 1:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            x + bw + Inches(0.02),
            y_pos + Inches(0.3) + bh / 2 - Inches(0.08),
            gp - Inches(0.04),
            Inches(0.16),
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = LIGHT_GREEN
        arrow.line.fill.background()

# Pre-project bar
pre_proj = "Pre-project:  Technical recommendations  |  RGPD actions  |  Accessibility planned  |  RICE prioritization"
ppbg = add_rounded_box(
    slide, Inches(0.5), Inches(4.6), SLIDE_WIDTH - Inches(1), Inches(0.55), RGBColor(0xE3, 0xF2, 0xFD), BLUE
)
add_textbox(
    slide,
    pre_proj,
    Inches(0.7),
    Inches(4.65),
    SLIDE_WIDTH - Inches(1.4),
    Inches(0.45),
    font_size=12,
    color=BLUE,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)
add_slide_number(slide, slide_num)

# ===================================================================
# SECTION: E2 - Architecture (C2, C3, C4, C6)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_slide(slide, "E2 -- Architecture", "Competencies C2, C3, C4, C6")
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 15: DATA TOPOGRAPHY (C2)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Data Topography -- 4 Parts", competency="C2")

topo_parts = [
    ("1", "Semantics", "9 business objects, metadata catalog", GREEN),
    ("2", "Data Models", "Structured (PG) + Semi (JSON) + Unstructured (Parquet)", BLUE),
    ("3", "Flows", "8 source-to-target flows, flux matrix, ETL diagrams", PURPLE),
    ("4", "Access", "Role-based matrix, 3 levels, RGPD constraints", ACCENT),
]

cols = 2
rows_t = 2
cw = Inches(5.5)
ch = Inches(1.4)
gx = Inches(0.5)
gy = Inches(0.4)
sx = Inches(0.6)
sy = Inches(1.8)

for idx, (num, title, desc, color) in enumerate(topo_parts):
    r = idx // cols
    c = idx % cols
    x = sx + c * (cw + gx)
    y = sy + r * (ch + gy)

    light = lighten((color[0], color[1], color[2]), 0.9)
    add_rounded_box(slide, x, y, cw, ch, light, color)

    # Number circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x + Inches(0.15), y + ch / 2 - Inches(0.25), Inches(0.5), Inches(0.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = num
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = WHITE

    add_textbox(
        slide,
        f"{title}: {desc}",
        x + Inches(0.8),
        y + Inches(0.2),
        cw - Inches(1.0),
        ch - Inches(0.4),
        font_size=13,
        color=DARK_GRAY,
        bold=False,
    )

add_why_callout(
    slide,
    "C2 requires mapping ALL data in 4 dimensions -- semantics, models, flows, and access conditions.",
    Inches(5.6),
)
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 16: SYSTEM ARCHITECTURE (C3)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "System Architecture -- 15 Services", competency="C3")

# Simplified architecture: rows of services
arch_rows = [
    ("SOURCES", GREEN, [("OFF API", None), ("Parquet", None), ("Scraping", None)]),
    ("ORCHESTRATION", PURPLE, [("Airflow (Web + Scheduler + Worker)", None)]),
    ("STORAGE", BLUE, [("PostgreSQL :5432", None), ("MinIO :9000", None), ("Redis :6379", None)]),
    ("SERVING", ACCENT, [("FastAPI :8000", "Devs"), ("Superset :8088", "Analysts"), ("Streamlit :8501", "Users")]),
    ("MONITORING", RED, [("Prometheus", None), ("Grafana :3000", None), ("MailHog :8025", None)]),
]

sy = Inches(1.5)
rh = Inches(0.7)
rg = Inches(0.35)

for ri, (row_label, color, services) in enumerate(arch_rows):
    y = sy + ri * (rh + rg)

    # Label
    lbl = add_rounded_box(slide, Inches(0.3), y, Inches(1.8), rh, color)
    add_textbox(
        slide,
        row_label,
        Inches(0.3),
        y + Inches(0.15),
        Inches(1.8),
        Inches(0.4),
        font_size=9,
        color=WHITE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # Services
    svc_x = Inches(2.3)
    n = len(services)
    svc_w = min(Inches(3.0), (SLIDE_WIDTH - Inches(2.8)) / n - Inches(0.1))
    for si, (svc_name, audience) in enumerate(services):
        x = svc_x + si * (svc_w + Inches(0.15))
        light = lighten((color[0], color[1], color[2]), 0.88)
        sbox = add_rounded_box(slide, x, y, svc_w, rh, light, color)
        add_textbox(
            slide,
            svc_name,
            x,
            y + Inches(0.1),
            svc_w,
            Inches(0.3),
            font_size=10,
            color=color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        if audience:
            add_textbox(
                slide,
                audience,
                x,
                y + Inches(0.4),
                svc_w,
                Inches(0.2),
                font_size=8,
                color=GRAY,
                bold=False,
                alignment=PP_ALIGN.CENTER,
            )

add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 17: FLUX MATRIX (C2, C3)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Flux Matrix -- 8 Data Flows", competency="C2 C3")

headers = ["Source", "Format", "Target", "Script", "Freq."]
rows_data = [
    ["OFF API", "JSON", "/data/raw/api/", "extract_off_api.py", "Daily"],
    ["OFF Parquet", "Parquet", "/data/raw/parquet/", "extract_off_parquet.py", "Weekly"],
    ["ANSES/EFSA", "HTML", "/data/raw/scraping/", "extract_scraping.py", "Monthly"],
    ["PostgreSQL", "SQL", "/data/raw/database/", "extract_from_db.py", "On-demand"],
    ["DuckDB", "Parquet", "/data/raw/duckdb/", "extract_duckdb.py", "On-demand"],
    ["Cleaned data", "Parquet", "PostgreSQL app", "import_to_db.py", "Daily"],
    ["App schema", "SQL", "DW star schema", "etl_load_warehouse", "Daily"],
    ["Raw files", "Mixed", "MinIO buckets", "etl_datalake_ingest", "Daily"],
]

add_table(slide, Inches(0.5), Inches(1.6), SLIDE_WIDTH - Inches(1), Inches(4.0), headers, rows_data)
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 18: TECHNICAL MONITORING (C4)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Technical Monitoring Newsletter", competency="C4")

topics = [
    ("Superset 6.0", "Migrated viz types\nLegacy to ECharts\n\nApplied", ACCENT),
    ("GDPR Update", "EDPB Guidelines 01/2026\nWellness = sensitive data\n\nCompliant", BLUE),
    ("Airflow 2.8", "Object storage backend\nListener hooks\n\nFuture", PURPLE),
]

card_w = Inches(3.5)
card_h = Inches(2.5)
sx = Inches(0.7)
gp = Inches(0.6)
y_pos = Inches(1.8)

for i, (title, desc, color) in enumerate(topics):
    x = sx + i * (card_w + gp)
    light = lighten((color[0], color[1], color[2]), 0.9)
    add_rounded_box(slide, x, y_pos, card_w, card_h, light, color)
    add_textbox(
        slide,
        title,
        x,
        y_pos + Inches(0.2),
        card_w,
        Inches(0.35),
        font_size=16,
        color=color,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        desc,
        x + Inches(0.2),
        y_pos + Inches(0.65),
        card_w - Inches(0.4),
        card_h - Inches(0.9),
        font_size=11,
        color=DARK_GRAY,
        bold=False,
        alignment=PP_ALIGN.CENTER,
    )

# Schedule bar
schedule_bar = add_rounded_box(
    slide, Inches(0.5), Inches(5.0), SLIDE_WIDTH - Inches(1), Inches(0.5), RGBColor(0xE8, 0xF5, 0xE9), GREEN
)
add_textbox(
    slide,
    "Weekly (1h/week)   |   Sources verified   |   docs/veille_technologique.md",
    Inches(0.7),
    Inches(5.05),
    SLIDE_WIDTH - Inches(1.4),
    Inches(0.4),
    font_size=12,
    color=GREEN,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 19: RGPD COMPLIANCE (C3, C11)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "RGPD Compliance -- How It Works", competency="C3 C11")

rgpd_cards = [
    ("Data Registry", "rgpd_data_registry table\nLegal basis per field\nRetention periods", RED),
    ("Consent", "consent_data_processing\nconsent_date column\nMandatory at registration", BLUE),
    ("Auto-Cleanup", "rgpd_cleanup_expired\nMeals >2y deleted\nUsers past retention", GREEN),
    ("Security", "bcrypt passwords\nUUID identification\nSHA256 in DW", PURPLE),
]

card_w = Inches(2.7)
card_h = Inches(2.2)
sx = Inches(0.4)
gp = Inches(0.3)
y_pos = Inches(1.8)

for i, (title, desc, color) in enumerate(rgpd_cards):
    x = sx + i * (card_w + gp)
    light = lighten((color[0], color[1], color[2]), 0.9)
    add_rounded_box(slide, x, y_pos, card_w, card_h, light, color)
    add_textbox(
        slide,
        title,
        x,
        y_pos + Inches(0.2),
        card_w,
        Inches(0.35),
        font_size=14,
        color=color,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        desc,
        x + Inches(0.15),
        y_pos + Inches(0.65),
        card_w - Inches(0.3),
        card_h - Inches(0.9),
        font_size=10,
        color=DARK_GRAY,
        bold=False,
        alignment=PP_ALIGN.CENTER,
    )

add_why_callout(
    slide,
    "Personal data stays in PostgreSQL (row-level delete). Public product data goes to both DW and Lake. This is the RGPD boundary.",
    Inches(4.6),
)
add_slide_number(slide, slide_num)

# ===================================================================
# SECTION: E3 - Project Kickoff (C5, C6, C7)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_slide(slide, "E3 -- Project Kickoff", "Competencies C5, C6, C7")
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 21: 6-PHASE ROADMAP (C5)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "6-Phase Roadmap", competency="C5")

phases = [
    ("Setup", "Docker\nDB init", "Wk 1-2"),
    ("Extract", "5 scripts\n3 DAGs", "Wk 3-4"),
    ("Transform", "Clean\nAggregate", "Wk 5-6"),
    ("Warehouse", "Star schema\nSCD", "Wk 7-8"),
    ("Lake", "Medallion\nCatalog", "Wk 9-10"),
    ("Deploy", "Monitor\nDocs", "Wk 11-12"),
]

sx = Inches(0.3)
bw = Inches(1.8)
bh = Inches(1.6)
gp = Inches(0.25)
y_pos = Inches(1.7)

for i, (title, desc, week) in enumerate(phases):
    x = sx + i * (bw + gp)
    light_green = RGBColor(0xE8, 0xF5, 0xE9)
    box = add_rounded_box(slide, x, y_pos, bw, bh, light_green, GREEN)
    add_textbox(
        slide,
        title,
        x,
        y_pos + Inches(0.15),
        bw,
        Inches(0.3),
        font_size=14,
        color=GREEN,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        desc,
        x,
        y_pos + Inches(0.5),
        bw,
        Inches(0.8),
        font_size=10,
        color=DARK_GRAY,
        bold=False,
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        week,
        x,
        y_pos + bh + Inches(0.05),
        bw,
        Inches(0.25),
        font_size=10,
        color=GRAY,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    if i < len(phases) - 1:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, x + bw + Inches(0.02), y_pos + bh / 2 - Inches(0.08), gp - Inches(0.04), Inches(0.16)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = LIGHT_GREEN
        arrow.line.fill.background()

# Bottom three info boxes
info_boxes = [
    ("Fibonacci story points", BLUE),
    ("Rituals: Planning, standup, review, retro", PURPLE),
    ("5 Roles: Data Eng, Platform, CI/CD, Auditor, Review", ACCENT),
]

ib_w = Inches(3.5)
ib_h = Inches(0.65)
ib_sx = Inches(0.8)
ib_gp = Inches(0.5)
ib_y = Inches(4.3)

for i, (text, color) in enumerate(info_boxes):
    x = ib_sx + i * (ib_w + ib_gp)
    light = lighten((color[0], color[1], color[2]), 0.9)
    add_rounded_box(slide, x, ib_y, ib_w, ib_h, light, color)
    add_textbox(
        slide,
        text,
        x + Inches(0.1),
        ib_y + Inches(0.1),
        ib_w - Inches(0.2),
        ib_h - Inches(0.2),
        font_size=10,
        color=color,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 22: TRACKING INDICATORS (C6)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Tracking Indicators", competency="C6")

metrics = [
    ("DAG Success Rate", "Pass/fail per run"),
    ("Activity Log", "22 entries, 14 days"),
    ("Completeness %", "Per dataset"),
    ("ETL Duration", "Trend over time"),
]

card_w = Inches(2.7)
card_h = Inches(1.8)
sx = Inches(0.5)
gp = Inches(0.35)
y_pos = Inches(2.0)

for i, (title, desc) in enumerate(metrics):
    x = sx + i * (card_w + gp)
    light_green = RGBColor(0xE8, 0xF5, 0xE9)
    add_rounded_box(slide, x, y_pos, card_w, card_h, light_green, GREEN)
    add_textbox(
        slide,
        title,
        x,
        y_pos + Inches(0.3),
        card_w,
        Inches(0.35),
        font_size=14,
        color=GREEN,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        desc,
        x,
        y_pos + Inches(0.8),
        card_w,
        Inches(0.5),
        font_size=11,
        color=DARK_GRAY,
        bold=False,
        alignment=PP_ALIGN.CENTER,
    )

# Tools bar
tools_bar = add_rounded_box(
    slide, Inches(0.5), Inches(4.5), SLIDE_WIDTH - Inches(1), Inches(0.5), RGBColor(0xE3, 0xF2, 0xFD), BLUE
)
add_textbox(
    slide,
    "Tools:  Airflow UI  |  Grafana SLA Dashboard  |  Git history",
    Inches(0.7),
    Inches(4.55),
    SLIDE_WIDTH - Inches(1.4),
    Inches(0.4),
    font_size=13,
    color=BLUE,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 23: MULTI-AUDIENCE COMMUNICATION (C7)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Multi-Audience Communication", competency="C7")

audiences = [
    ("Developers", ["Swagger / ReDoc", "Grafana dashboards", "Git + MkDocs"], GREEN),
    ("Analysts", ["Superset dashboards", "DW datamart views", "Data catalog"], BLUE),
    ("End Users", ["Streamlit app", "Product search", "Meal tracking"], ACCENT),
]

card_w = Inches(3.5)
card_h = Inches(2.5)
sx = Inches(0.7)
gp = Inches(0.6)
y_pos = Inches(1.8)

for i, (title, items, color) in enumerate(audiences):
    x = sx + i * (card_w + gp)
    light = lighten((color[0], color[1], color[2]), 0.9)
    add_rounded_box(slide, x, y_pos, card_w, card_h, light, color)
    add_textbox(
        slide,
        title,
        x,
        y_pos + Inches(0.2),
        card_w,
        Inches(0.35),
        font_size=16,
        color=color,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_bullet_list(
        slide,
        items,
        x + Inches(0.3),
        y_pos + Inches(0.65),
        card_w - Inches(0.6),
        card_h - Inches(0.9),
        font_size=12,
        color=DARK_GRAY,
        bullet_color=color,
    )

add_why_callout(
    slide,
    "Right tool for right audience. Technical stakeholders get Swagger; business gets Superset; users get Streamlit.",
    Inches(5.0),
)
add_slide_number(slide, slide_num)

# ===================================================================
# SECTION: E4 - Data Collection & API (C8-C12)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_slide(slide, "E4 -- Data Collection & API", "Competencies C8, C9, C10, C11, C12")
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 25: 5 EXTRACTION SOURCES (C8)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "5 Extraction Sources", competency="C8")

sources = [
    ("REST API", "1k+ products/day", GREEN),
    ("Data File", "50k Parquet", BLUE),
    ("Scraping", "ANSES/EFSA", PURPLE),
    ("Database", "PostgreSQL", ACCENT),
    ("Big Data", "DuckDB 3M+", RED),
]

src_w = Inches(2.0)
src_h = Inches(1.4)
sx = Inches(0.3)
gp = Inches(0.3)
y_pos = Inches(1.8)

for i, (title, desc, color) in enumerate(sources):
    x = sx + i * (src_w + gp)
    light = lighten((color[0], color[1], color[2]), 0.88)
    add_rounded_box(slide, x, y_pos, src_w, src_h, light, color)
    add_textbox(
        slide,
        title,
        x,
        y_pos + Inches(0.2),
        src_w,
        Inches(0.3),
        font_size=14,
        color=color,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        desc,
        x,
        y_pos + Inches(0.6),
        src_w,
        Inches(0.4),
        font_size=10,
        color=DARK_GRAY,
        bold=False,
        alignment=PP_ALIGN.CENTER,
    )

# Arrow down to aggregate_clean
for i in range(5):
    x = sx + i * (src_w + gp) + src_w / 2 - Inches(0.08)
    darrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, y_pos + src_h + Inches(0.05), Inches(0.16), Inches(0.3))
    darrow.fill.solid()
    darrow.fill.fore_color.rgb = LIGHT_GRAY
    darrow.line.fill.background()

# aggregate_clean box
agg_w = Inches(4.0)
agg_x = (SLIDE_WIDTH - agg_w) / 2
agg_y = y_pos + src_h + Inches(0.5)
agg_box = add_rounded_box(slide, agg_x, agg_y, agg_w, Inches(0.6), RGBColor(0xE8, 0xF5, 0xE9), GREEN)
add_textbox(
    slide,
    "aggregate_clean.py",
    agg_x,
    agg_y + Inches(0.1),
    agg_w,
    Inches(0.4),
    font_size=14,
    color=GREEN,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)

add_why_callout(
    slide, "C8 requires 5 source types: REST API, data file, web scraping, database, big data system.", Inches(4.8)
)
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 26: REST API DETAIL (C8)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Source 1: Open Food Facts REST API", competency="C8")

# Left column - steps
steps_left = [
    "Configure 5 food categories",
    "Paginated GET requests",
    "Rate limiting (0.6s delay)",
    "User-Agent header",
    "Try/except per page",
    "Save JSON to /data/raw/api/",
]

add_textbox(
    slide, "extract_off_api.py", Inches(0.5), Inches(1.5), Inches(5), Inches(0.35), font_size=16, color=GREEN, bold=True
)
add_bullet_list(
    slide,
    steps_left,
    Inches(0.5),
    Inches(2.0),
    Inches(5),
    Inches(3.5),
    font_size=13,
    color=DARK_GRAY,
    bullet_color=GREEN,
)

# Right column - script structure
struct_items = [
    "1. Entry point: main()",
    "2. Dependency init",
    "3. External connections",
    "4. Logic rules",
    "5. Error handling",
    "6. Result saving",
    "7. Versioned on Git",
]

light_blue = RGBColor(0xE3, 0xF2, 0xFD)
add_rounded_box(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(3.5), light_blue, BLUE)
add_textbox(
    slide,
    "Script structure (all 5):",
    Inches(7.0),
    Inches(1.65),
    Inches(5),
    Inches(0.3),
    font_size=14,
    color=BLUE,
    bold=True,
)
add_bullet_list(
    slide,
    struct_items,
    Inches(7.0),
    Inches(2.1),
    Inches(5),
    Inches(2.8),
    font_size=12,
    color=DARK_GRAY,
    bullet_color=BLUE,
)

add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 27: SCRAPING + DUCKDB (C8)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Source 3: Web Scraping & Source 5: DuckDB", competency="C8")

# Left: Scraping
scraping_items = [
    "Targets: ANSES + EFSA",
    "BeautifulSoup parsing",
    "Nutritional guidelines",
    "Fallback RDA values",
    "JSON output",
    "Frequency: Monthly",
]

light_purple = RGBColor(0xF3, 0xE5, 0xF5)
add_rounded_box(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(4.0), light_purple, PURPLE)
add_textbox(
    slide,
    "Web Scraping",
    Inches(0.5),
    Inches(1.6),
    Inches(5.5),
    Inches(0.35),
    font_size=18,
    color=PURPLE,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)
add_textbox(
    slide,
    "extract_scraping.py",
    Inches(0.5),
    Inches(2.0),
    Inches(5.5),
    Inches(0.3),
    font_size=13,
    color=PURPLE,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)
add_bullet_list(
    slide,
    scraping_items,
    Inches(0.8),
    Inches(2.5),
    Inches(5),
    Inches(2.5),
    font_size=12,
    color=DARK_GRAY,
    bullet_color=PURPLE,
)

# Right: DuckDB
duckdb_items = [
    "SQL on Parquet files",
    "3M+ rows, no DB needed",
    "Analytical aggregations",
    "Columnar = fast scans",
    "Parquet output",
    "Think: lightweight Spark",
]

light_red = RGBColor(0xFC, 0xE4, 0xE4)
add_rounded_box(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(4.0), light_red, RED)
add_textbox(
    slide,
    "DuckDB Big Data",
    Inches(6.8),
    Inches(1.6),
    Inches(5.5),
    Inches(0.35),
    font_size=18,
    color=RED,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)
add_textbox(
    slide,
    "extract_duckdb.py",
    Inches(6.8),
    Inches(2.0),
    Inches(5.5),
    Inches(0.3),
    font_size=13,
    color=RED,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)
add_bullet_list(
    slide,
    duckdb_items,
    Inches(7.1),
    Inches(2.5),
    Inches(5),
    Inches(2.5),
    font_size=12,
    color=DARK_GRAY,
    bullet_color=RED,
)

add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 28: CLEANING PIPELINE (C10)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Cleaning Pipeline", competency="C10")

add_textbox(
    slide,
    "aggregate_clean.py -- 7 cleaning steps",
    Inches(0.5),
    Inches(1.4),
    Inches(10),
    Inches(0.4),
    font_size=16,
    color=DARK_GRAY,
    bold=True,
)

cleaning_steps = [
    ("1.", "Standardize 30+ column names across all sources", RED),
    ("2.", "Validate barcodes (strip non-numeric, check 8-14 digits)", ACCENT),
    ("3.", "Remove rows with null product names", ACCENT),
    ("4.", "Cap nutritional values at physiological max per 100g", PURPLE),
    ("5.", "Normalize Nutri-Score to uppercase A-E", BLUE),
    ("6.", "Deduplicate by barcode (keep most complete)", GREEN),
    ("7.", "Generate cleaning_report.json", GREEN),
]

sy = Inches(2.0)
rh = Inches(0.55)
rg = Inches(0.1)

for i, (num, desc, color) in enumerate(cleaning_steps):
    y = sy + i * (rh + rg)
    light = lighten((color[0], color[1], color[2]), 0.92)
    add_rounded_box(slide, Inches(0.5), y, SLIDE_WIDTH - Inches(1), rh, light, color)
    add_textbox(
        slide,
        f"{num}  {desc}",
        Inches(0.7),
        y + Inches(0.08),
        SLIDE_WIDTH - Inches(1.4),
        Inches(0.4),
        font_size=13,
        color=DARK_GRAY,
        bold=False,
    )

add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 29: SQL QUERIES (C9)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "7 Optimized SQL Queries", competency="C9")

sql_headers = ["#", "Query", "Technique", "Optimization"]
sql_rows = [
    ["1", "Full-text product search", "GIN index, ts_rank", "Index scan"],
    ["2", "User nutrition summary", "ROW_NUMBER() OVER", "Window partition"],
    ["3", "Product market analysis", "CTE + aggregation", "CTE avoids re-scan"],
    ["4", "Brand comparison", "Analytical functions", "Composite index"],
    ["5", "Temporal trends", "LAG(), moving avg", "Partial index"],
    ["6", "Category aggregation", "GROUP BY + HAVING", "Aggregate pushdown"],
    ["7", "Complex joins", "Materialized + joins", "Pre-computed"],
]

add_table(slide, Inches(0.5), Inches(1.6), SLIDE_WIDTH - Inches(1), Inches(3.5), sql_headers, sql_rows)

add_textbox(
    slide,
    "All queries documented with EXPLAIN ANALYZE and optimization rationale.",
    Inches(0.5),
    Inches(5.5),
    SLIDE_WIDTH - Inches(1),
    Inches(0.4),
    font_size=13,
    color=GRAY,
    bold=False,
)
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 30: DATABASE (C11)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "RGPD-Compliant Database", competency="C11")

# Left: table diagram
tables = [
    ("users", BLUE, Inches(0.8), Inches(2.0)),
    ("meals", BLUE, Inches(0.8), Inches(2.7)),
    ("meal_items", BLUE, Inches(0.8), Inches(3.4)),
    ("countries", BLUE, Inches(0.8), Inches(4.1)),
    ("products", BLUE, Inches(3.5), Inches(2.0)),
    ("brands", BLUE, Inches(3.5), Inches(2.7)),
    ("categories", BLUE, Inches(3.5), Inches(3.4)),
    ("rgpd_registry", RED, Inches(3.5), Inches(4.1)),
]

for name, color, x, y in tables:
    light = lighten((color[0], color[1], color[2]), 0.9)
    add_rounded_box(slide, x, y, Inches(2.0), Inches(0.5), light, color)
    add_textbox(
        slide,
        name,
        x,
        y + Inches(0.08),
        Inches(2.0),
        Inches(0.35),
        font_size=11,
        color=color,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

# Right: details
db_details = [
    "MERISE: MCD -> MLD -> MPD",
    "rgpd_data_registry table",
    "Consent columns",
    "rgpd_cleanup_expired_data()",
    "bcrypt password hashing",
    "UUID-based identification",
    "8 tables, full referential integrity",
]

light_green = RGBColor(0xE8, 0xF5, 0xE9)
add_rounded_box(slide, Inches(6.5), Inches(1.8), Inches(5.8), Inches(3.8), light_green, GREEN)
add_bullet_list(
    slide,
    db_details,
    Inches(6.8),
    Inches(2.0),
    Inches(5.2),
    Inches(3.4),
    font_size=13,
    color=DARK_GRAY,
    bullet_color=GREEN,
)
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 31: FASTAPI REST API (C12)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "FastAPI REST API", competency="C12")

# Left: Endpoints
endpoints = [
    "POST /auth/register",
    "POST /auth/login",
    "GET /products/?q=&nutriscore=",
    "GET /products/{barcode}",
    "GET /products/{barcode}/alternatives",
    "POST /meals",
    "GET /meals/daily-summary",
    "GET /meals/weekly-trends",
]

add_textbox(
    slide, "Endpoints", Inches(0.5), Inches(1.5), Inches(5), Inches(0.35), font_size=16, color=ACCENT, bold=True
)

y_ep = Inches(2.0)
for i, ep in enumerate(endpoints):
    light_accent = RGBColor(0xFF, 0xF3, 0xE0)
    add_rounded_box(slide, Inches(0.5), y_ep + i * Inches(0.45), Inches(5.2), Inches(0.38), light_accent, ACCENT)
    add_textbox(
        slide,
        ep,
        Inches(0.7),
        y_ep + i * Inches(0.45) + Inches(0.04),
        Inches(4.8),
        Inches(0.3),
        font_size=11,
        color=DARK_GRAY,
        bold=False,
        font_name="Courier New",
    )

# Right: Auth flow
add_textbox(slide, "Auth Flow", Inches(7), Inches(1.5), Inches(5), Inches(0.35), font_size=16, color=BLUE, bold=True)

auth_steps = [
    ("Login", "username + password", ACCENT),
    ("JWT Token", "60-min expiry", BLUE),
    ("Role Check", "user / nutritionist / admin", GREEN),
    ("Response", "Pydantic validated", PURPLE),
]

auth_y = Inches(2.0)
auth_w = Inches(4.5)

for i, (title, desc, color) in enumerate(auth_steps):
    y = auth_y + i * Inches(1.0)
    light = lighten((color[0], color[1], color[2]), 0.9)
    add_rounded_box(slide, Inches(7.5), y, auth_w, Inches(0.7), light, color)
    add_textbox(
        slide,
        title,
        Inches(7.5),
        y + Inches(0.05),
        auth_w,
        Inches(0.3),
        font_size=12,
        color=color,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        desc,
        Inches(7.5),
        y + Inches(0.35),
        auth_w,
        Inches(0.25),
        font_size=9,
        color=GRAY,
        bold=False,
        alignment=PP_ALIGN.CENTER,
    )

    if i < len(auth_steps) - 1:
        darrow = slide.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW, Inches(7.5) + auth_w / 2 - Inches(0.08), y + Inches(0.72), Inches(0.16), Inches(0.2)
        )
        darrow.fill.solid()
        darrow.fill.fore_color.rgb = LIGHT_GRAY
        darrow.line.fill.background()

add_demo_callout(slide, "Show Swagger UI at localhost:8000/docs", Inches(6.2))
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 32: STREAMLIT - TWO ROLES (C7, C20)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Streamlit -- Two Roles", competency="C7 C20")

# User App
user_items = [
    "Product search + filter",
    "Log meals",
    "Daily macro dashboard",
    "Weekly trends",
    "Product comparison",
    "Healthier alternatives",
    "Consumes FastAPI REST API",
]

light_accent = RGBColor(0xFF, 0xF3, 0xE0)
add_rounded_box(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(4.0), light_accent, ACCENT)
add_textbox(
    slide,
    "User App (port 8501)",
    Inches(0.5),
    Inches(1.6),
    Inches(5.5),
    Inches(0.35),
    font_size=18,
    color=ACCENT,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)
add_bullet_list(
    slide,
    user_items,
    Inches(0.8),
    Inches(2.1),
    Inches(5),
    Inches(3.0),
    font_size=12,
    color=DARK_GRAY,
    bullet_color=ACCENT,
)

# Data Catalog
catalog_items = [
    "Search datasets",
    "Browse bronze/silver/gold",
    "View schema + lineage",
    "Quality metrics",
    "Governance & RGPD info",
    "Storage statistics",
    "Connects directly to MinIO",
]

light_purple = RGBColor(0xF3, 0xE5, 0xF5)
add_rounded_box(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(4.0), light_purple, PURPLE)
add_textbox(
    slide,
    "Data Catalog (C20)",
    Inches(6.8),
    Inches(1.6),
    Inches(5.5),
    Inches(0.35),
    font_size=18,
    color=PURPLE,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)
add_bullet_list(
    slide,
    catalog_items,
    Inches(7.1),
    Inches(2.1),
    Inches(5),
    Inches(3.0),
    font_size=12,
    color=DARK_GRAY,
    bullet_color=PURPLE,
)

add_demo_callout(slide, "Show both Streamlit pages at localhost:8501", Inches(6.0))
add_slide_number(slide, slide_num)

# ===================================================================
# SECTION: E5 - Data Warehouse (C13-C15)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_slide(slide, "E5 -- Data Warehouse", "Competencies C13, C14, C15")
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 34: STAR SCHEMA (C13)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Star Schema", competency="C13")

# Fact tables (center)
facts = [
    ("fact_daily\nnutrition", ACCENT, Inches(3.5), Inches(3.0)),
    ("fact_product\nmarket", ACCENT, Inches(8.0), Inches(3.0)),
]

dims = [
    ("dim_time", BLUE, Inches(1.0), Inches(1.6)),
    ("dim_user", BLUE, Inches(3.5), Inches(1.6)),
    ("dim_product", BLUE, Inches(6.0), Inches(1.6)),
    ("dim_nutriscore", BLUE, Inches(1.0), Inches(4.8)),
    ("dim_brand", BLUE, Inches(5.5), Inches(4.8)),
    ("dim_category", BLUE, Inches(8.0), Inches(4.8)),
    ("dim_country", BLUE, Inches(10.5), Inches(4.8)),
]

# Draw facts
for name, color, x, y in facts:
    box = add_rounded_box(slide, x, y, Inches(2.2), Inches(1.0), RGBColor(0xFF, 0xF3, 0xE0), ACCENT)
    add_textbox(
        slide,
        name,
        x,
        y + Inches(0.15),
        Inches(2.2),
        Inches(0.7),
        font_size=12,
        color=ACCENT,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

# Draw dims
for name, color, x, y in dims:
    light_blue_bg = RGBColor(0xE3, 0xF2, 0xFD)
    box = add_rounded_box(slide, x, y, Inches(1.8), Inches(0.65), light_blue_bg, BLUE)
    add_textbox(
        slide,
        name,
        x,
        y + Inches(0.1),
        Inches(1.8),
        Inches(0.4),
        font_size=10,
        color=BLUE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

# Bottom summary
summary_bar = add_rounded_box(
    slide, Inches(1.5), Inches(6.0), SLIDE_WIDTH - Inches(3), Inches(0.5), RGBColor(0xE8, 0xF5, 0xE9), GREEN
)
add_textbox(
    slide,
    "7 dims  +  2 facts  +  6 datamart views    |    Bottom-up (Kimball)",
    Inches(1.5),
    Inches(6.05),
    SLIDE_WIDTH - Inches(3),
    Inches(0.4),
    font_size=13,
    color=GREEN,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 35: 6 DATAMART VIEWS (C13, C14)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "6 Datamart Views", competency="C13 C14")

dm_headers = ["Datamart View", "Purpose", "Audience"]
dm_rows = [
    ["dm_user_daily_nutrition", "Daily calorie & macro tracking", "Nutritionists"],
    ["dm_product_market_by_category", "Product availability by category", "Market analysts"],
    ["dm_brand_quality_ranking", "Brand ranking by Nutri-Score", "Quality team"],
    ["dm_nutriscore_distribution", "Nutri-Score distribution (A-E)", "Public health"],
    ["dm_nutrition_trends", "Weekly/monthly nutrition trends", "Data scientists"],
    ["dm_dw_health", "DW operational health metrics", "DW admins"],
]

add_table(slide, Inches(0.5), Inches(1.6), SLIDE_WIDTH - Inches(1), Inches(3.2), dm_headers, dm_rows)

add_why_callout(
    slide,
    "Datamarts are pre-built SQL views so analysts don't need to write complex star-schema joins themselves.",
    Inches(5.3),
)
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 36: ETL PIPELINE - 7 DAGS (C15)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "ETL Pipeline -- 7 DAGs", competency="C15")

# Extract DAGs (column 1)
extract_dags = [
    ("extract_off_api", GREEN),
    ("extract_parquet", GREEN),
    ("extract_scraping", GREEN),
]
add_textbox(slide, "02:00", Inches(0.3), Inches(2.8), Inches(0.8), Inches(0.3), font_size=10, color=GRAY, bold=True)
for i, (name, color) in enumerate(extract_dags):
    y = Inches(1.7) + i * Inches(0.65)
    light = lighten((color[0], color[1], color[2]), 0.88)
    add_rounded_box(slide, Inches(1.0), y, Inches(2.5), Inches(0.5), light, color)
    add_textbox(
        slide,
        name,
        Inches(1.0),
        y + Inches(0.08),
        Inches(2.5),
        Inches(0.35),
        font_size=10,
        color=color,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

# Arrows to clean
for i in range(3):
    y_src = Inches(1.7) + i * Inches(0.65) + Inches(0.25)
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.6), y_src - Inches(0.08), Inches(0.6), Inches(0.16))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = LIGHT_GRAY
    arrow.line.fill.background()

# Clean DAG (column 2)
add_textbox(slide, "04:00", Inches(4.3), Inches(3.0), Inches(0.8), Inches(0.3), font_size=10, color=GRAY, bold=True)
light_accent = lighten((ACCENT[0], ACCENT[1], ACCENT[2]), 0.88)
add_rounded_box(slide, Inches(4.3), Inches(2.2), Inches(2.5), Inches(0.5), light_accent, ACCENT)
add_textbox(
    slide,
    "aggregate_clean",
    Inches(4.3),
    Inches(2.28),
    Inches(2.5),
    Inches(0.35),
    font_size=10,
    color=ACCENT,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)

# Arrows to DW and Lake
arrow_r1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.9), Inches(2.15), Inches(0.6), Inches(0.16))
arrow_r1.fill.solid()
arrow_r1.fill.fore_color.rgb = LIGHT_GRAY
arrow_r1.line.fill.background()

arrow_r2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.9), Inches(2.75), Inches(0.6), Inches(0.16))
arrow_r2.fill.solid()
arrow_r2.fill.fore_color.rgb = LIGHT_GRAY
arrow_r2.line.fill.background()

# DW + Lake DAGs (column 3)
add_textbox(slide, "05:00", Inches(7.6), Inches(3.5), Inches(0.8), Inches(0.3), font_size=10, color=GRAY, bold=True)
add_textbox(
    slide,
    "parallel",
    Inches(7.6),
    Inches(1.7),
    Inches(2.5),
    Inches(0.3),
    font_size=10,
    color=PURPLE,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)

light_blue = lighten((BLUE[0], BLUE[1], BLUE[2]), 0.88)
add_rounded_box(slide, Inches(7.6), Inches(2.0), Inches(2.5), Inches(0.5), light_blue, BLUE)
add_textbox(
    slide,
    "load_warehouse",
    Inches(7.6),
    Inches(2.08),
    Inches(2.5),
    Inches(0.35),
    font_size=10,
    color=BLUE,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)

light_purple = lighten((PURPLE[0], PURPLE[1], PURPLE[2]), 0.88)
add_rounded_box(slide, Inches(7.6), Inches(2.6), Inches(2.5), Inches(0.5), light_purple, PURPLE)
add_textbox(
    slide,
    "datalake_ingest",
    Inches(7.6),
    Inches(2.68),
    Inches(2.5),
    Inches(0.35),
    font_size=10,
    color=PURPLE,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)

# Backup DAG
add_textbox(slide, "06:00", Inches(0.3), Inches(4.5), Inches(0.8), Inches(0.3), font_size=10, color=GRAY, bold=True)
light_red = lighten((RED[0], RED[1], RED[2]), 0.88)
add_rounded_box(slide, Inches(1.0), Inches(4.3), Inches(2.5), Inches(0.5), light_red, RED)
add_textbox(
    slide,
    "backup_maint.",
    Inches(1.0),
    Inches(4.38),
    Inches(2.5),
    Inches(0.35),
    font_size=10,
    color=RED,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)

# Sensor note
sensor_bg = add_rounded_box(slide, Inches(6.5), Inches(4.0), Inches(5.5), Inches(0.7), RGBColor(0xE3, 0xF2, 0xFD), BLUE)
add_textbox(
    slide,
    "ExternalTaskSensor ensures DW and Lake wait for cleaning to finish",
    Inches(6.7),
    Inches(4.1),
    Inches(5.1),
    Inches(0.5),
    font_size=11,
    color=BLUE,
    bold=False,
)

add_demo_callout(slide, "Show Airflow UI at localhost:8080", Inches(5.3))
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 37: SCD (C17)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "SCD -- Slowly Changing Dimensions", competency="C17")

scds = [
    (
        "Type 1 -- Overwrite",
        "(dim_brand)",
        "Brand name misspelled? Just overwrite it. No history needed for corrections.",
        GREEN,
    ),
    (
        "Type 2 -- Historical",
        "(dim_product)",
        "Product changes? Close old row (end_date, is_current=false), insert new row. Full history.",
        BLUE,
    ),
    (
        "Type 3 -- Previous Value",
        "(dim_country)",
        "Country list changes? Keep previous_country_list column. One level of history.",
        PURPLE,
    ),
]

sy = Inches(1.8)
rh = Inches(1.3)
rg = Inches(0.3)

for i, (title, table_name, desc, color) in enumerate(scds):
    y = sy + i * (rh + rg)
    light = lighten((color[0], color[1], color[2]), 0.92)
    add_rounded_box(slide, Inches(0.5), y, SLIDE_WIDTH - Inches(1), rh, light, color)
    add_textbox(
        slide,
        f"{title}  {table_name}",
        Inches(0.8),
        y + Inches(0.15),
        SLIDE_WIDTH - Inches(1.6),
        Inches(0.35),
        font_size=16,
        color=color,
        bold=True,
    )
    add_textbox(
        slide,
        desc,
        Inches(0.8),
        y + Inches(0.55),
        SLIDE_WIDTH - Inches(1.6),
        Inches(0.6),
        font_size=12,
        color=DARK_GRAY,
        bold=False,
    )

add_why_callout(
    slide,
    "Change detection uses IS DISTINCT FROM in the ETL. All 3 types are integrated into the warehouse loading DAG.",
    Inches(6.2),
)
add_slide_number(slide, slide_num)

# ===================================================================
# SECTION: E6 - DW Maintenance (C16, C17)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_slide(slide, "E6 -- DW Maintenance", "Competencies C16, C17", bg_color=BLUE)
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 39: ALERT SYSTEM (C16)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Alert System -- How It Works", competency="C16")

alert_steps = [
    ("DAG Fails", RED),
    ("alerting.py\ncallback", ACCENT),
    ("activity_log", BLUE),
    ("MailHog", PURPLE),
    ("StatsD", GREEN),
    ("Grafana", GREEN),
]

sx = Inches(0.3)
bw = Inches(1.8)
bh = Inches(1.1)
gp = Inches(0.25)
y_pos = Inches(2.5)

for i, (label, color) in enumerate(alert_steps):
    x = sx + i * (bw + gp)
    light = lighten((color[0], color[1], color[2]), 0.88)
    add_rounded_box(slide, x, y_pos, bw, bh, light, color)
    add_textbox(
        slide,
        label,
        x,
        y_pos + Inches(0.2),
        bw,
        bh - Inches(0.4),
        font_size=11,
        color=color,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    if i < len(alert_steps) - 1:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, x + bw + Inches(0.02), y_pos + bh / 2 - Inches(0.08), gp - Inches(0.04), Inches(0.16)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = LIGHT_GRAY
        arrow.line.fill.background()

add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 40: SLA DASHBOARD (C16)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "SLA Dashboard", competency="C16")

sla_metrics = [
    ("ETL Success", "> 95%"),
    ("Freshness", "< 24h"),
    ("Backups", "100%"),
    ("Query Time", "< 5s"),
]

card_w = Inches(2.7)
card_h = Inches(1.6)
sx = Inches(0.5)
gp = Inches(0.35)
y_pos = Inches(1.8)

for i, (title, target) in enumerate(sla_metrics):
    x = sx + i * (card_w + gp)
    light_green = RGBColor(0xE8, 0xF5, 0xE9)
    add_rounded_box(slide, x, y_pos, card_w, card_h, light_green, GREEN)
    add_textbox(
        slide,
        title,
        x,
        y_pos + Inches(0.2),
        card_w,
        Inches(0.35),
        font_size=14,
        color=GREEN,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        target,
        x,
        y_pos + Inches(0.7),
        card_w,
        Inches(0.4),
        font_size=18,
        color=GREEN,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

# ITIL Priority
itil_bg = add_rounded_box(
    slide, Inches(0.5), Inches(4.0), SLIDE_WIDTH - Inches(1), Inches(0.9), RGBColor(0xFF, 0xF3, 0xE0), ACCENT
)
add_textbox(
    slide,
    "ITIL Priority Matrix:  P1 Critical <1h  |  P2 High <4h  |  P3 Medium <24h  |  P4 Low: next sprint\nEscalation: L1 (auto-restart) -> L2 (engineer) -> L3 (architecture review)",
    Inches(0.7),
    Inches(4.05),
    SLIDE_WIDTH - Inches(1.4),
    Inches(0.8),
    font_size=11,
    color=ACCENT,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)

add_demo_callout(slide, "Show Grafana SLA dashboard + MailHog alerts", Inches(5.3))
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 41: BACKUP & MAINTENANCE (C16)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Backup & Maintenance Procedures", competency="C16")

# Left: Backup Script
backup_items = [
    "backup_database.py",
    "Full + partial modes",
    "Upload to MinIO /backups",
    "Auto-cleanup old backups",
    "Airflow DAG scheduled",
    "154 lines, Git-versioned",
]

light_blue = RGBColor(0xE3, 0xF2, 0xFD)
add_rounded_box(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(4.0), light_blue, BLUE)
add_textbox(
    slide,
    "Backup Script",
    Inches(0.5),
    Inches(1.6),
    Inches(5.5),
    Inches(0.35),
    font_size=18,
    color=BLUE,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)
add_bullet_list(
    slide,
    backup_items,
    Inches(0.8),
    Inches(2.1),
    Inches(5),
    Inches(3.0),
    font_size=13,
    color=DARK_GRAY,
    bullet_color=BLUE,
)

# Right: Documented Procedures
proc_items = [
    "New source integration (6 steps)",
    "New access creation (3 steps)",
    "Storage expansion",
    "Add datamart view",
    "Compute capacity scaling",
    "All in rapport_final Ch.12",
]

light_green = RGBColor(0xE8, 0xF5, 0xE9)
add_rounded_box(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(4.0), light_green, GREEN)
add_textbox(
    slide,
    "Documented Procedures",
    Inches(6.8),
    Inches(1.6),
    Inches(5.5),
    Inches(0.35),
    font_size=18,
    color=GREEN,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)
add_bullet_list(
    slide,
    proc_items,
    Inches(7.1),
    Inches(2.1),
    Inches(5),
    Inches(3.0),
    font_size=13,
    color=DARK_GRAY,
    bullet_color=GREEN,
)

add_slide_number(slide, slide_num)

# ===================================================================
# SECTION: E7 - Data Lake (C18-C21)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_slide(slide, "E7 -- Data Lake", "Competencies C18, C19, C20, C21", bg_color=PURPLE)
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 43: MEDALLION ARCHITECTURE (C18)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Medallion Architecture", competency="C18")

layers = [
    ("Bronze", "Raw as-is\nJSON, Parquet, CSV\n90-day retention", BRONZE),
    ("Silver", "Cleaned & validated\nParquet only\nDeduplicated", SILVER),
    ("Gold", "Analytics-ready\nML features\nDaily snapshots", GOLD),
]

card_w = Inches(3.3)
card_h = Inches(2.5)
sx = Inches(0.6)
gp = Inches(0.8)
y_pos = Inches(1.8)

for i, (title, desc, color) in enumerate(layers):
    x = sx + i * (card_w + gp)
    light = lighten((color[0], color[1], color[2]), 0.85)
    box = add_rounded_box(slide, x, y_pos, card_w, card_h, light, color)
    add_textbox(
        slide,
        title,
        x,
        y_pos + Inches(0.2),
        card_w,
        Inches(0.4),
        font_size=22,
        color=color,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        desc,
        x + Inches(0.2),
        y_pos + Inches(0.8),
        card_w - Inches(0.4),
        card_h - Inches(1.0),
        font_size=12,
        color=DARK_GRAY,
        bold=False,
        alignment=PP_ALIGN.CENTER,
    )

    if i < len(layers) - 1:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            x + card_w + Inches(0.1),
            y_pos + card_h / 2 - Inches(0.12),
            gp - Inches(0.2),
            Inches(0.24),
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = color
        arrow.line.fill.background()
        label_text = "clean" if i == 0 else "aggregate"
        add_textbox(
            slide,
            label_text,
            x + card_w + Inches(0.1),
            y_pos + card_h / 2 - Inches(0.35),
            gp - Inches(0.2),
            Inches(0.2),
            font_size=10,
            color=color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

# V/V/V bar
vvv_bar = add_rounded_box(
    slide, Inches(0.5), Inches(4.9), SLIDE_WIDTH - Inches(1), Inches(0.5), RGBColor(0xE8, 0xF5, 0xE9), GREEN
)
add_textbox(
    slide,
    "Addresses: Volume (3M+ products)  |  Variety (JSON, Parquet, CSV)  |  Velocity (daily/weekly schedules)",
    Inches(0.7),
    Inches(4.95),
    SLIDE_WIDTH - Inches(1.4),
    Inches(0.4),
    font_size=12,
    color=GREEN,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 44: VOLUME, VARIETY, VELOCITY (C18)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Volume, Variety, Velocity", competency="C18")

vvv = [
    ("Volume", ["3M+ products", "50k+ imported", "Parquet columnar", "for fast scans"], BLUE),
    ("Variety", ["JSON (API)", "Parquet (bulk)", "HTML (scraping)", "SQL (database)"], PURPLE),
    ("Velocity", ["Daily API pulls", "Weekly Parquet", "Monthly scraping", "Scheduled DAGs"], ACCENT),
]

card_w = Inches(3.5)
card_h = Inches(3.5)
sx = Inches(0.7)
gp = Inches(0.6)
y_pos = Inches(1.8)

for i, (title, items, color) in enumerate(vvv):
    x = sx + i * (card_w + gp)
    light = lighten((color[0], color[1], color[2]), 0.9)
    add_rounded_box(slide, x, y_pos, card_w, card_h, light, color)
    add_textbox(
        slide,
        title,
        x,
        y_pos + Inches(0.2),
        card_w,
        Inches(0.45),
        font_size=22,
        color=color,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_bullet_list(
        slide,
        items,
        x + Inches(0.3),
        y_pos + Inches(0.8),
        card_w - Inches(0.6),
        card_h - Inches(1.0),
        font_size=13,
        color=DARK_GRAY,
        bullet_color=color,
    )

add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 45: CATALOG COMPARISON (C18)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Catalog Tool Comparison", competency="C18")

cat_headers = ["Criteria", "Apache Atlas", "DataHub", "Custom JSON"]
cat_rows = [
    ["Setup complexity", "High", "Medium", "Low"],
    ["Resource overhead", "Heavy (Java)", "Moderate", "Minimal"],
    ["MinIO integration", "Manual", "Plugin", "Native"],
    ["Search capability", "Full", "Full", "Basic (Streamlit)"],
    ["Lineage tracking", "Auto", "Auto", "Manual (JSON)"],
    ["Fits project scale", "Overkill", "Overkill", "Right-sized"],
]

add_table(slide, Inches(0.5), Inches(1.6), SLIDE_WIDTH - Inches(1), Inches(3.2), cat_headers, cat_rows)

add_why_callout(
    slide,
    "Custom JSON catalog was chosen because Atlas and DataHub add significant infrastructure for a project at this scale. The Streamlit browser provides search and browse capabilities.",
    Inches(5.3),
)
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 46: DATA CATALOG BROWSER (C20)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Data Catalog Browser", competency="C20")

add_textbox(
    slide,
    "Streamlit Catalog Page",
    Inches(0.5),
    Inches(1.5),
    Inches(10),
    Inches(0.4),
    font_size=18,
    color=PURPLE,
    bold=True,
)

catalog_features = [
    "Search datasets by name, format, or source",
    "Browse by layer (bronze / silver / gold tabs)",
    "View schema (column names + types)",
    "View lineage (where data came from)",
    "Quality metrics per dataset",
    "Governance info + RGPD compliance",
    "Storage stats (object count, total size)",
]

sy = Inches(2.1)
rh = Inches(0.48)
rg = Inches(0.1)

for i, feat in enumerate(catalog_features):
    y = sy + i * (rh + rg)
    light_purple = RGBColor(0xF3, 0xE5, 0xF5)
    add_rounded_box(slide, Inches(0.5), y, Inches(8.5), rh, light_purple, PURPLE)
    add_textbox(
        slide, feat, Inches(0.7), y + Inches(0.06), Inches(8.1), Inches(0.35), font_size=12, color=DARK_GRAY, bold=False
    )

add_demo_callout(slide, "Show Streamlit catalog at localhost:8501/data_catalog", Inches(6.2))
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 47: ACCESS GOVERNANCE (C21)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Access Governance", competency="C21")

add_textbox(
    slide,
    "Group-Based Access (not individual)",
    Inches(0.5),
    Inches(1.5),
    Inches(10),
    Inches(0.4),
    font_size=18,
    color=DARK_GRAY,
    bold=True,
)

roles = [
    ("admin_role", "PostgreSQL: full  |  MinIO: all buckets  |  API: admin  |  Superset: Admin", GREEN),
    (
        "nutritionist_role",
        "PostgreSQL: read products+meals  |  MinIO: gold  |  API: nutritionist  |  Superset: Analyst",
        BLUE,
    ),
    (
        "app_readonly",
        "PostgreSQL: read products only  |  MinIO: gold (public)  |  API: user  |  Superset: Viewer",
        ACCENT,
    ),
]

sy = Inches(2.2)
rh = Inches(0.75)
rg = Inches(0.2)

for i, (role, perms, color) in enumerate(roles):
    y = sy + i * (rh + rg)
    light = lighten((color[0], color[1], color[2]), 0.92)
    add_rounded_box(slide, Inches(0.5), y, SLIDE_WIDTH - Inches(1), rh, light, color)

    # Role name
    txBox = slide.shapes.add_textbox(Inches(0.7), y + Inches(0.1), Inches(2.5), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = role
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = BODY_FONT

    # Permissions
    add_textbox(
        slide, perms, Inches(3.3), y + Inches(0.15), Inches(9), Inches(0.45), font_size=11, color=DARK_GRAY, bold=False
    )

add_why_callout(
    slide,
    "Least-privilege principle: each role gets only what it needs. Consistent across all 4 systems. RGPD-compliant.",
    Inches(5.2),
)
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 48: MONITORING STACK (C16, C20)
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Monitoring Stack", competency="C16 C20")

# Top row: 4 exporters
exporters = [
    ("StatsD", "Airflow metrics"),
    ("cAdvisor", "Container stats"),
    ("Node Exp.", "Host metrics"),
    ("PG Exp.", "DB metrics"),
]

exp_w = Inches(2.5)
exp_h = Inches(1.2)
sx_exp = Inches(0.5)
gp_exp = Inches(0.35)
y_exp = Inches(1.7)

for i, (name, desc) in enumerate(exporters):
    x = sx_exp + i * (exp_w + gp_exp)
    light_red = RGBColor(0xFC, 0xE4, 0xE4)
    add_rounded_box(slide, x, y_exp, exp_w, exp_h, light_red, RED)
    add_textbox(
        slide,
        name,
        x,
        y_exp + Inches(0.15),
        exp_w,
        Inches(0.3),
        font_size=13,
        color=RED,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        desc,
        x,
        y_exp + Inches(0.5),
        exp_w,
        Inches(0.3),
        font_size=10,
        color=DARK_GRAY,
        bold=False,
        alignment=PP_ALIGN.CENTER,
    )

    # Down arrow
    darrow = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, x + exp_w / 2 - Inches(0.08), y_exp + exp_h + Inches(0.05), Inches(0.16), Inches(0.3)
    )
    darrow.fill.solid()
    darrow.fill.fore_color.rgb = LIGHT_GRAY
    darrow.line.fill.background()

# Prometheus
prom_y = y_exp + exp_h + Inches(0.5)
light_accent = RGBColor(0xFF, 0xF3, 0xE0)
add_rounded_box(slide, Inches(2.0), prom_y, Inches(8), Inches(0.7), light_accent, ACCENT)
add_textbox(
    slide,
    "Prometheus   (scrapes all 4 exporters)",
    Inches(2.0),
    prom_y + Inches(0.1),
    Inches(8),
    Inches(0.5),
    font_size=14,
    color=ACCENT,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)

# Down arrow to Grafana
darrow2 = slide.shapes.add_shape(
    MSO_SHAPE.DOWN_ARROW, Inches(6) - Inches(0.1), prom_y + Inches(0.75), Inches(0.2), Inches(0.3)
)
darrow2.fill.solid()
darrow2.fill.fore_color.rgb = LIGHT_GRAY
darrow2.line.fill.background()

# Grafana
graf_y = prom_y + Inches(1.2)
light_green = RGBColor(0xE8, 0xF5, 0xE9)
add_rounded_box(slide, Inches(3.0), graf_y, Inches(6), Inches(0.7), light_green, GREEN)
add_textbox(
    slide,
    "Grafana   (6 dashboards)",
    Inches(3.0),
    graf_y + Inches(0.1),
    Inches(6),
    Inches(0.5),
    font_size=14,
    color=GREEN,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)

add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 49: 6 GRAFANA DASHBOARDS
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "6 Grafana Dashboards")

graf_headers = ["Dashboard", "Monitors", "Key Panels"]
graf_rows = [
    ["Airflow Overview", "DAG runs, task durations", "Success rate, active tasks"],
    ["Airflow DAGs", "Per-DAG performance", "Run outcomes, SLA misses"],
    ["PostgreSQL", "DB health, connections", "Query time, rows, cache"],
    ["Docker System", "Container resources", "CPU, memory, network"],
    ["MinIO", "Object storage", "Bucket sizes, requests"],
    ["SLA Compliance", "Service indicators", "ETL %, freshness, backups"],
]

add_table(slide, Inches(0.5), Inches(1.6), SLIDE_WIDTH - Inches(1), Inches(3.2), graf_headers, graf_rows)

add_demo_callout(slide, "Show Grafana at localhost:3000", Inches(5.3))
add_slide_number(slide, slide_num)

# ===================================================================
# SECTION: CONCLUSION
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_slide(slide, "Conclusion", None, bg_color=GREEN)
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 51: 21/21 COMPETENCIES
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "21/21 Competencies Covered")

blocks_comp = [
    ("Block 1: Steer", GREEN, ["C1", "C2", "C3", "C4", "C5", "C6", "C7"]),
    ("Block 2: Collect", BLUE, ["C8", "C9", "C10", "C11", "C12"]),
    ("Block 3: Warehouse", PURPLE, ["C13", "C14", "C15", "C16", "C17"]),
    ("Block 4: Lake", ACCENT, ["C18", "C19", "C20", "C21"]),
]

sy = Inches(1.6)
rg = Inches(0.7)

for bi, (block_name, color, comps) in enumerate(blocks_comp):
    y = sy + bi * rg

    # Block label
    lbl = add_rounded_box(slide, Inches(0.5), y, Inches(2.5), Inches(0.5), color)
    add_textbox(
        slide,
        block_name,
        Inches(0.5),
        y + Inches(0.08),
        Inches(2.5),
        Inches(0.35),
        font_size=12,
        color=WHITE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # Competency boxes
    cx = Inches(3.3)
    cw = Inches(0.7)
    for ci, comp in enumerate(comps):
        x = cx + ci * (cw + Inches(0.1))
        light = lighten((color[0], color[1], color[2]), 0.85)
        cbox = add_rounded_box(slide, x, y + Inches(0.05), cw, Inches(0.4), light, color)
        add_textbox(
            slide,
            comp,
            x,
            y + Inches(0.1),
            cw,
            Inches(0.3),
            font_size=9,
            color=color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

# 21/21 badge
add_textbox(
    slide,
    "21/21",
    Inches(10),
    Inches(2.0),
    Inches(2),
    Inches(0.6),
    font_size=36,
    color=GREEN,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)

# Summary bar
summary_bar = add_rounded_box(
    slide, Inches(0.5), Inches(4.8), SLIDE_WIDTH - Inches(1), Inches(0.55), RGBColor(0xE8, 0xF5, 0xE9), GREEN
)
add_textbox(
    slide,
    "One integrated project  |  15 Docker services  |  Full RGPD compliance  |  All evaluations E1-E7",
    Inches(0.7),
    Inches(4.85),
    SLIDE_WIDTH - Inches(1.4),
    Inches(0.45),
    font_size=13,
    color=GREEN,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)
add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 52: LESSONS LEARNED
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Lessons Learned")

# Left: What worked
worked = [
    "Docker: 1-command deploy",
    "DW + Lake in parallel",
    "Airflow orchestration",
    "RGPD as architecture driver",
    "Docs as first-class output",
]

light_green = RGBColor(0xE8, 0xF5, 0xE9)
add_rounded_box(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(4.0), light_green, GREEN)
add_textbox(
    slide,
    "What worked",
    Inches(0.5),
    Inches(1.6),
    Inches(5.5),
    Inches(0.4),
    font_size=18,
    color=GREEN,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)
add_bullet_list(
    slide, worked, Inches(0.8), Inches(2.2), Inches(5), Inches(3.0), font_size=14, color=DARK_GRAY, bullet_color=GREEN
)

# Right: Next steps
next_steps = [
    "Real-time (Kafka)",
    "Kubernetes",
    "ML pipeline",
    "DataHub catalog",
    "More tests",
]

light_accent = RGBColor(0xFF, 0xF3, 0xE0)
add_rounded_box(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(4.0), light_accent, ACCENT)
add_textbox(
    slide,
    "Next steps",
    Inches(6.8),
    Inches(1.6),
    Inches(5.5),
    Inches(0.4),
    font_size=18,
    color=ACCENT,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)
add_bullet_list(
    slide,
    next_steps,
    Inches(7.1),
    Inches(2.2),
    Inches(5),
    Inches(3.0),
    font_size=14,
    color=DARK_GRAY,
    bullet_color=ACCENT,
)

add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 53: LIVE DEMO PLAN
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Live Demo Plan")

add_textbox(
    slide,
    "Demo Sequence (during the 60-minute presentation)",
    Inches(0.5),
    Inches(1.4),
    Inches(10),
    Inches(0.4),
    font_size=16,
    color=DARK_GRAY,
    bold=True,
)

demos = [
    ("1.", "docker compose up -d -- show all 15 services starting", GREEN),
    ("2.", "Swagger UI (localhost:8000/docs) -- register, login, search products", ACCENT),
    ("3.", "Streamlit (localhost:8501) -- search product, log meal, see dashboard", BLUE),
    ("4.", "Airflow UI (localhost:8080) -- show DAG graph, trigger a run", PURPLE),
    ("5.", "Superset (localhost:8088) -- show BI dashboard", BLUE),
    ("6.", "Grafana (localhost:3000) -- show SLA dashboard", RED),
    ("7.", "MailHog (localhost:8025) -- show alert emails", PURPLE),
    ("8.", "MinIO Console (localhost:9001) -- browse bronze/silver/gold buckets", ACCENT),
    ("9.", "Streamlit Catalog (localhost:8501/data_catalog) -- search datasets", PURPLE),
]

sy = Inches(1.9)
rh = Inches(0.48)
rg = Inches(0.08)

for i, (num, desc, color) in enumerate(demos):
    y = sy + i * (rh + rg)
    light = lighten((color[0], color[1], color[2]), 0.92)
    add_rounded_box(slide, Inches(0.5), y, SLIDE_WIDTH - Inches(1), rh, light, color)
    add_textbox(
        slide,
        f"{num}  {desc}",
        Inches(0.7),
        y + Inches(0.06),
        SLIDE_WIDTH - Inches(1.4),
        Inches(0.35),
        font_size=12,
        color=DARK_GRAY,
        bold=False,
    )

add_slide_number(slide, slide_num)

# ===================================================================
# SLIDE 54: THANK YOU
# ===================================================================
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, GREEN)

# NT logo
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.666), Inches(1.0), Inches(2.0), Inches(2.0))
circle.fill.solid()
circle.fill.fore_color.rgb = WHITE
circle.line.fill.background()
tf = circle.text_frame
tf.word_wrap = False
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "NT"
run.font.size = Pt(48)
run.font.bold = True
run.font.color.rgb = GREEN
run.font.name = TITLE_FONT

add_textbox(
    slide,
    "Thank You",
    Inches(1),
    Inches(3.3),
    SLIDE_WIDTH - Inches(2),
    Inches(1),
    font_size=54,
    color=WHITE,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)
add_textbox(
    slide,
    "Questions?",
    Inches(1),
    Inches(4.3),
    SLIDE_WIDTH - Inches(2),
    Inches(0.6),
    font_size=28,
    color=RGBColor(0xC8, 0xE6, 0xC9),
    bold=False,
    alignment=PP_ALIGN.CENTER,
)
add_textbox(
    slide,
    "github.com/Reetika12795/NutriTrack",
    Inches(1),
    Inches(5.5),
    SLIDE_WIDTH - Inches(2),
    Inches(0.4),
    font_size=16,
    color=WHITE,
    bold=False,
    alignment=PP_ALIGN.CENTER,
)
add_textbox(
    slide,
    "reetika12795.github.io/NutriTrack",
    Inches(1),
    Inches(5.9),
    SLIDE_WIDTH - Inches(2),
    Inches(0.4),
    font_size=16,
    color=WHITE,
    bold=False,
    alignment=PP_ALIGN.CENTER,
)
add_slide_number(slide, slide_num)

# ===================================================================
# SAVE
# ===================================================================
output_path = os.path.join(os.path.dirname(__file__), "NutriTrack_Defense_Visual.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {slide_num}")
